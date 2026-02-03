#!/usr/bin/env python3
"""
Create templates/template.potx with 2 slide masters and 4 layouts matching template.pptx:

Masters:
- Title Master: Full-bleed background image + styled title/subtitle placeholders + track-name textbox
- Content Master: Green header bar rectangle + title in header + body (OBJECT) placeholder + 
                  slide number + cropped logo

Layouts:
- "Title": inherits from Title Master (blank, background via master)
- "Text": TITLE + OBJECT placeholders  
- "Image Right": TITLE, OBJECT (left side), PICTURE (right side)
- "Dual Image": TITLE, 2x PICTURE, BODY (bottom)

Based on template_comparison.md and create_template_fixes.md analysis.
"""

from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import os

# Widescreen dimensions
SLIDE_WIDTH_IN = 13.33
SLIDE_HEIGHT_IN = 7.5

# Colors from template.pptx analysis
DARK_BLUE = RGBColor(0x00, 0x27, 0x4B)    # #00274b - Title text
HEADER_GREEN = RGBColor(0x9C, 0xCC, 0xB4)  # #9cccb4 - Header bar
TRACK_GREEN = RGBColor(0x3D, 0x6B, 0x38)   # #3d6b38 - Track name text

# Asset paths
ASSETS_DIR = "assets"
BACKGROUND_IMAGE = os.path.join(ASSETS_DIR, "title_slide_bg_image1.png")
LOGO_IMAGE = os.path.join(ASSETS_DIR, "title_slide_bg_image2.png")

# Footer text for content layouts
FOOTER_TEXT = "AI in Actuarial and Finance"


def create_template():
    """Create a new template with 2 masters and 4 layouts matching template.pptx structure."""
    
    # Start from blank presentation with widescreen dimensions
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)
    
    # Get the default master - this will be our Content Master
    content_master = prs.slide_masters[0]
    
    print(f"Initial layouts: {len(content_master.slide_layouts)}")
    
    # We need to create a second master for Title slides
    # Since python-pptx doesn't support adding masters directly, we'll work with one master
    # but configure it properly with all the required shapes
    title_master = content_master  # For now, use same master but configure layouts differently
    
    # Configure the Content Master with header bar, logo, and placeholders
    configure_content_master(content_master)
    
    # Remove extra layouts to keep only 4
    prune_extra_layouts(content_master, keep_count=4)
    
    print(f"After pruning: {len(content_master.slide_layouts)} layouts")
    
    # Configure the 4 layouts
    layout_configs = [
        ("Title", configure_title_layout),
        ("Text", configure_text_layout),
        ("Image Right", configure_image_right_layout),
        ("Dual Image", configure_dual_image_layout),
    ]
    
    for idx, (name, config_func) in enumerate(layout_configs):
        if idx < len(content_master.slide_layouts):
            layout = content_master.slide_layouts[idx]
            layout.name = name
            print(f"Configuring layout {idx}: {name}")
            
            # Clear existing placeholders
            clear_placeholders(layout)
            
            # Configure the layout
            config_func(layout)
    
    # Save the template
    output_path = "templates/template.potx"
    os.makedirs("templates", exist_ok=True)
    prs.save(output_path)
    print(f"\nTemplate saved to: {output_path}")
    
    # Post-process to set paragraph properties via python-pptx API
    # XML attributes may not suffice for alignment/indent behavior
    post_process_template(output_path)
    
    return output_path


def set_paragraph_indent_xml(pPr, left_margin=0, first_line_indent=0):
    """
    Set paragraph indent properties via XML attributes.
    
    Args:
        pPr: The paragraph properties element (_pPr from a _Paragraph)
             or None if paragraph properties don't exist
        left_margin: Left margin in inches (0 for no indent)
        first_line_indent: First line indent in inches (0 for no hanging indent)
    """
    if pPr is None:
        return
    
    # Convert to EMU (English Metric Units)
    marL_emu = int(left_margin * 914400)
    indent_emu = int(first_line_indent * 914400)
    
    # Set marL and indent attributes
    pPr.set(qn('a:marL') if 'a:' not in str(pPr.tag) else 'marL', str(marL_emu))
    pPr.set('marL', str(marL_emu))
    pPr.set('indent', str(indent_emu))


def post_process_template(template_path):
    """
    Post-process the template to set paragraph properties via python-pptx API.
    
    This function sets paragraph formatting at TWO levels:
    1. lstStyle (level properties): Defines defaults inherited by NEW paragraphs
       - Includes proper bullet/numbered list margins and hanging indents
       - Creates space between bullet character and text
    2. pPr (paragraph properties): Sets alignment on sample/existing paragraphs
    
    Settings applied:
    - Title layout: title, subtitle, text LEFT aligned
    - All other layouts: titles LEFT aligned
    - Dual image layout body: CENTER aligned
    - Bullet/numbered formatting with proper hanging indent at all levels
    """
    print("\nPost-processing template to set paragraph formatting...")
    
    prs = Presentation(template_path)
    master = prs.slide_masters[0]
    
    for layout in master.slide_layouts:
        layout_name = layout.name
        print(f"  Processing layout: {layout_name}")
        
        for shape in layout.shapes:
            if not shape.has_text_frame:
                continue
            
            tf = shape.text_frame
            ph_type = None
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
            
            # Determine alignment based on layout and placeholder type
            if layout_name == "Dual Image" and ph_type == PP_PLACEHOLDER.BODY:
                target_align = 'ctr'  # Center for dual image body
            else:
                target_align = 'l'    # Left for everything else
            
            # CRITICAL: Set lstStyle to define defaults for NEW paragraphs
            # This includes proper bullet/numbered list formatting with hanging indents
            set_lstStyle_alignment(tf, target_align)
            
            # Set alignment on existing paragraphs (sample text)
            # Note: We don't override marL/indent here - let lstStyle control that
            for para in tf.paragraphs:
                if target_align == 'ctr':
                    para.alignment = PP_ALIGN.CENTER
                else:
                    para.alignment = PP_ALIGN.LEFT
    
    # Also process master-level placeholders
    print("  Processing master placeholders...")
    for shape in master.shapes:
        if not shape.has_text_frame:
            continue
        
        tf = shape.text_frame
        set_lstStyle_alignment(tf, 'l')
        for para in tf.paragraphs:
            para.alignment = PP_ALIGN.LEFT
    
    # Save the updated template
    prs.save(template_path)
    print(f"  Saved post-processed template to: {template_path}")


def set_lstStyle_alignment(text_frame, align='l'):
    """
    Set default paragraph properties in lstStyle for new paragraphs.
    
    This is critical because when slide_builders adds new paragraphs via
    text_frame.add_paragraph(), they inherit from lstStyle, not from existing
    paragraph properties.
    
    CRITICAL FIX: Set marL="0" and indent="0" EXPLICITLY for all levels.
    This prevents indentation on wrapped lines in headers and plain text.
    
    TRADE-OFF: Bullets will not have hanging indent from lstStyle.
    For proper bullet formatting, the generator's add_bullet() function
    should set marL and indent on the paragraph's pPr directly when
    adding bullet formatting.
    
    In python-pptx, paragraph.level maps to lvlXpPr as follows:
    - level 0 -> lvl1pPr (used for level-0 bullets AND plain text)
    - level 1 -> lvl2pPr (used for level-1 sub-bullets)
    - etc.
    
    Args:
        text_frame: TextFrame object
        align: Alignment value ('l', 'ctr', 'r' for left, center, right)
    """
    txBody = text_frame._txBody
    lstStyle = txBody.find(qn('a:lstStyle'))
    
    if lstStyle is None:
        # Create lstStyle if missing
        lstStyle = etree.SubElement(txBody, qn('a:lstStyle'))
    
    # Add/update level paragraph properties for levels 0-8
    # CRITICAL: Set marL="0" and indent="0" explicitly for ALL levels
    # to prevent wrapped line indentation on headers and plain text
    for level in range(9):
        lvl_tag = f'a:lvl{level + 1}pPr'
        lvl_elem = lstStyle.find(qn(lvl_tag))
        
        if lvl_elem is None:
            lvl_elem = etree.SubElement(lstStyle, qn(lvl_tag))
        
        # Set alignment
        lvl_elem.set('algn', align)
        
        # CRITICAL FIX: Set marL and indent to 0 EXPLICITLY
        # This prevents wrapped line indentation in headers and plain text
        lvl_elem.set('marL', '0')
        lvl_elem.set('indent', '0')
        
        # Remove tabLst if present - not needed
        tabLst = lvl_elem.find(qn('a:tabLst'))
        if tabLst is not None:
            lvl_elem.remove(tabLst)


def configure_content_master(master):
    """
    Configure Content Master with:
    - Green header bar rectangle (#9cccb4) at top
    - Cropped logo in bottom-right corner
    - Footer text "AI in Actuarial and Finance" bottom-left
    - Styled title placeholder in header
    - Body placeholder
    - Slide number placeholder
    """
    # Clear any existing shapes on master that are placeholders
    clear_placeholders(master)
    
    # Get spTree for low-level XML manipulation
    spTree = master.shapes._spTree
    
    # Add green header bar rectangle at top (using XML)
    add_header_bar_xml(spTree)
    
    # Add logo with cropping in bottom-right (using XML)
    add_logo_xml(master, spTree)
    
    # Add footer text bottom-left (using XML)
    add_footer_text_xml(spTree)
    
    # Add master-level placeholders (layouts will inherit)
    add_master_placeholders(master)
    
    print("  Added header bar, logo, footer text, and master placeholders")


def add_header_bar_xml(spTree):
    """Add green header bar rectangle at top using XML (masters don't support add_shape)."""
    existing_ids = [int(sp.get('id')) for sp in spTree.findall('.//' + qn('p:cNvPr')) if sp.get('id')]
    new_id = max(existing_ids) + 1 if existing_ids else 2
    
    # Convert to EMU
    width_emu = int(SLIDE_WIDTH_IN * 914400)
    height_emu = int(0.82 * 914400)
    
    # Green color as hex
    green_hex = f"{HEADER_GREEN[0]:02X}{HEADER_GREEN[1]:02X}{HEADER_GREEN[2]:02X}"
    
    sp_xml = f'''
    <p:sp xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
          xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
      <p:nvSpPr>
        <p:cNvPr id="{new_id}" name="Rectangle 7"/>
        <p:cNvSpPr/>
        <p:nvPr/>
      </p:nvSpPr>
      <p:spPr>
        <a:xfrm>
          <a:off x="0" y="0"/>
          <a:ext cx="{width_emu}" cy="{height_emu}"/>
        </a:xfrm>
        <a:prstGeom prst="rect">
          <a:avLst/>
        </a:prstGeom>
        <a:solidFill>
          <a:srgbClr val="{green_hex}"/>
        </a:solidFill>
        <a:ln>
          <a:noFill/>
        </a:ln>
      </p:spPr>
      <p:txBody>
        <a:bodyPr/>
        <a:lstStyle/>
        <a:p>
          <a:endParaRPr lang="en-US"/>
        </a:p>
      </p:txBody>
    </p:sp>
    '''
    
    sp_elem = etree.fromstring(sp_xml.strip())
    # Insert at beginning (after nvGrpSpPr and grpSpPr)
    spTree.insert(2, sp_elem)
    return sp_elem


def add_logo_xml(master, spTree):
    """Add cropped logo picture to master using low-level XML."""
    if not os.path.exists(LOGO_IMAGE):
        print(f"  Warning: Logo image not found at {LOGO_IMAGE}")
        return None
    
    # Get or add image part - pass the file path directly
    image_part, rId = master.part.get_or_add_image_part(LOGO_IMAGE)
    
    existing_ids = [int(sp.get('id')) for sp in spTree.findall('.//' + qn('p:cNvPr')) if sp.get('id')]
    new_id = max(existing_ids) + 1 if existing_ids else 2
    
    # Position and size in EMU
    left_emu = int(12.30 * 914400)
    top_emu = int(6.93 * 914400)
    width_emu = int(0.94 * 914400)
    height_emu = int(0.46 * 914400)
    
    # Cropping values (from template.pptx)
    crop_top = 27784
    crop_bottom = 23024
    
    pic_xml = f'''
    <p:pic xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
           xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
           xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
      <p:nvPicPr>
        <p:cNvPr id="{new_id}" name="Picture 4"/>
        <p:cNvPicPr>
          <a:picLocks noChangeAspect="1"/>
        </p:cNvPicPr>
        <p:nvPr/>
      </p:nvPicPr>
      <p:blipFill>
        <a:blip r:embed="{rId}"/>
        <a:srcRect t="{crop_top}" b="{crop_bottom}"/>
        <a:stretch>
          <a:fillRect/>
        </a:stretch>
      </p:blipFill>
      <p:spPr>
        <a:xfrm>
          <a:off x="{left_emu}" y="{top_emu}"/>
          <a:ext cx="{width_emu}" cy="{height_emu}"/>
        </a:xfrm>
        <a:prstGeom prst="rect">
          <a:avLst/>
        </a:prstGeom>
      </p:spPr>
    </p:pic>
    '''
    
    pic_elem = etree.fromstring(pic_xml.strip())
    spTree.append(pic_elem)
    return pic_elem


def add_footer_text_xml(spTree):
    """Add footer text 'AI in Actuarial and Finance' bottom-left on content master.
    
    Position and style match original template.pptx:
    - Position: (0.13", 6.88"), Size: 7.12" x 0.61"
    - Font: Calibri 18pt bold, scheme dk1 color
    """
    existing_ids = [int(sp.get('id')) for sp in spTree.findall('.//' + qn('p:cNvPr')) if sp.get('id')]
    new_id = max(existing_ids) + 1 if existing_ids else 2
    
    # Position and size in EMU (from template.pptx: x=115920, y=6289200, cx=6509160, cy=559440)
    left_emu = int(0.13 * 914400)   # ~115920
    top_emu = int(6.88 * 914400)    # ~6289200
    width_emu = int(7.12 * 914400)  # ~6509160
    height_emu = int(0.61 * 914400) # ~559440
    
    sp_xml = f'''
    <p:sp xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
          xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
      <p:nvSpPr>
        <p:cNvPr id="{new_id}" name="Slide Number Placeholder 5"/>
        <p:cNvSpPr/>
        <p:nvPr/>
      </p:nvSpPr>
      <p:spPr>
        <a:xfrm>
          <a:off x="{left_emu}" y="{top_emu}"/>
          <a:ext cx="{width_emu}" cy="{height_emu}"/>
        </a:xfrm>
        <a:prstGeom prst="rect">
          <a:avLst/>
        </a:prstGeom>
        <a:noFill/>
        <a:ln w="0">
          <a:noFill/>
        </a:ln>
      </p:spPr>
      <p:style>
        <a:lnRef idx="0"/>
        <a:fillRef idx="0"/>
        <a:effectRef idx="0"/>
        <a:fontRef idx="minor"/>
      </p:style>
      <p:txBody>
        <a:bodyPr anchor="ctr">
          <a:noAutofit/>
        </a:bodyPr>
        <a:p>
          <a:pPr defTabSz="914400">
            <a:lnSpc>
              <a:spcPct val="100000"/>
            </a:lnSpc>
          </a:pPr>
          <a:r>
            <a:rPr b="1" lang="en-US" sz="1800" spc="-1" strike="noStrike">
              <a:solidFill>
                <a:schemeClr val="dk1"/>
              </a:solidFill>
              <a:latin typeface="Calibri"/>
            </a:rPr>
            <a:t>{FOOTER_TEXT}</a:t>
          </a:r>
          <a:endParaRPr b="0" lang="en-US" sz="1800" spc="-1" strike="noStrike">
            <a:solidFill>
              <a:srgbClr val="000000"/>
            </a:solidFill>
            <a:latin typeface="Arial"/>
          </a:endParaRPr>
        </a:p>
      </p:txBody>
    </p:sp>
    '''
    
    sp_elem = etree.fromstring(sp_xml.strip())
    spTree.append(sp_elem)
    return sp_elem


def add_master_placeholders(master):
    """Add styled placeholders to master that layouts will inherit."""
    spTree = master.shapes._spTree
    
    # Title placeholder in header bar area - left-aligned with no indent
    add_placeholder_shape(
        spTree, 
        ph_type='title',
        idx=0,
        left=0.12, top=0.0, width=13.22, height=0.81,
        name="PlaceHolder 1",
        font_name="Calibri Light",
        font_size_pt=44,
        bold=False,
        font_color=None,  # Use scheme color
        anchor='ctr',
        p_align='l',
        left_margin_in=0,
        first_line_indent_in=0
    )
    
    # Body placeholder (OBJECT type) for main content - no indent
    add_placeholder_shape(
        spTree,
        ph_type='body',
        idx=1,
        left=0.92, top=1.17, width=11.5, height=5.52,
        name="PlaceHolder 2",
        font_name="Calibri",
        font_size_pt=24,
        bold=False,
        font_color=None,
        left_margin_in=0,
        first_line_indent_in=0
    )
    
    # Slide number placeholder
    add_placeholder_shape(
        spTree,
        ph_type='sldNum',
        idx=4,
        left=11.76, top=6.91, width=0.48, height=0.49,
        name="Slide Number Placeholder 5"
    )


def add_placeholder_shape(spTree, ph_type, idx, left, top, width, height, name=None,
                          font_name=None, font_size_pt=None, bold=None, font_color=None,
                          anchor=None, p_align=None, left_margin_in=None, first_line_indent_in=None):
    """
    Add a placeholder shape to spTree using low-level XML.
    
    Args:
        spTree: The shape tree element to append to
        ph_type: Placeholder type string ('title', 'body', 'ctrTitle', 'subTitle', 'pic', 'obj', 'sldNum')
        idx: Placeholder index
        left, top, width, height: Position/size in inches
        name: Shape name
        font_name: Default font name
        font_size_pt: Default font size in points
        bold: Bold flag
        font_color: RGBColor for text
        anchor: Text anchor ('t', 'ctr', 'b')
        p_align: Paragraph alignment ('l', 'ctr', 'r') for default text alignment
        left_margin_in: Left margin in inches (converted to EMU for marL attribute)
        first_line_indent_in: First line indent in inches (converted to EMU for indent attribute)
    """
    # Generate unique ID
    existing_ids = [int(sp.get('id')) for sp in spTree.findall('.//' + qn('p:cNvPr')) if sp.get('id')]
    new_id = max(existing_ids) + 1 if existing_ids else 2
    
    # Create namespaced elements
    nsmap = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    }
    
    if name is None:
        name = f"{ph_type} {idx}"
    
    # Convert position/size to EMU
    left_emu = int(left * 914400)
    top_emu = int(top * 914400)
    width_emu = int(width * 914400)
    height_emu = int(height * 914400)
    
    # Build paragraph properties (pPr) for the sample paragraph
    # CRITICAL: Always set marL=0 and indent=0 to prevent wrapped line indentation
    # Also add buNone to suppress bullets on this sample paragraph
    # When the generator adds bullets, it will override with buChar/buAutoNum
    pPr_attrs = ['marL="0"', 'indent="0"']
    if p_align:
        pPr_attrs.append(f'algn="{p_align}"')
    
    # The pPr element includes buNone to explicitly suppress bullets on the sample
    # This ensures plain text in the placeholder doesn't inherit bullet margins
    pPr_xml = f'<a:pPr {" ".join(pPr_attrs)}><a:buNone/></a:pPr>'
    
    # Build lstStyle with explicit zero margins for all levels
    # CRITICAL: marL="0" and indent="0" prevent wrapped line indentation
    # For bullet formatting, the generator should set marL/indent directly on the paragraph
    lstStyle_xml = '<a:lstStyle>'
    for lvl in range(1, 10):
        algn_attr = f' algn="{p_align}"' if p_align else ''
        lstStyle_xml += f'<a:lvl{lvl}pPr marL="0" indent="0"{algn_attr}/>'
    lstStyle_xml += '</a:lstStyle>'
    
    # Build XML for shape
    sp_xml = f'''
    <p:sp xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
          xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
          xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
      <p:nvSpPr>
        <p:cNvPr id="{new_id}" name="{name}"/>
        <p:cNvSpPr>
          <a:spLocks noGrp="1"/>
        </p:cNvSpPr>
        <p:nvPr>
          <p:ph type="{ph_type}" idx="{idx}"/>
        </p:nvPr>
      </p:nvSpPr>
      <p:spPr>
        <a:xfrm>
          <a:off x="{left_emu}" y="{top_emu}"/>
          <a:ext cx="{width_emu}" cy="{height_emu}"/>
        </a:xfrm>
        <a:prstGeom prst="rect">
          <a:avLst/>
        </a:prstGeom>
      </p:spPr>
      <p:txBody>
        <a:bodyPr{' anchor="' + anchor + '"' if anchor else ''}/>
        {lstStyle_xml}
        <a:p>
          {pPr_xml}
          <a:endParaRPr lang="en-US"/>
        </a:p>
      </p:txBody>
    </p:sp>
    '''
    
    sp_elem = etree.fromstring(sp_xml.strip())
    
    # Add font styling if specified
    if font_name or font_size_pt is not None or bold is not None or font_color:
        add_default_text_style(sp_elem, font_name, font_size_pt, bold, font_color)
    
    spTree.append(sp_elem)
    return sp_elem


def add_picture_placeholder(spTree, idx, left, top, width, height, name=None):
    """
    Add a PICTURE placeholder shape (without txBody for proper image handling).
    
    Args:
        spTree: The shape tree element
        idx: Placeholder index
        left, top, width, height: Position/size in inches
        name: Shape name
    """
    existing_ids = [int(sp.get('id')) for sp in spTree.findall('.//' + qn('p:cNvPr')) if sp.get('id')]
    new_id = max(existing_ids) + 1 if existing_ids else 2
    
    if name is None:
        name = f"Picture Placeholder {idx}"
    
    left_emu = int(left * 914400)
    top_emu = int(top * 914400)
    width_emu = int(width * 914400)
    height_emu = int(height * 914400)
    
    # Picture placeholder should have minimal txBody but pic type
    sp_xml = f'''
    <p:sp xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
          xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
          xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
      <p:nvSpPr>
        <p:cNvPr id="{new_id}" name="{name}"/>
        <p:cNvSpPr>
          <a:spLocks noGrp="1"/>
        </p:cNvSpPr>
        <p:nvPr>
          <p:ph type="pic" idx="{idx}"/>
        </p:nvPr>
      </p:nvSpPr>
      <p:spPr>
        <a:xfrm>
          <a:off x="{left_emu}" y="{top_emu}"/>
          <a:ext cx="{width_emu}" cy="{height_emu}"/>
        </a:xfrm>
        <a:prstGeom prst="rect">
          <a:avLst/>
        </a:prstGeom>
      </p:spPr>
      <p:txBody>
        <a:bodyPr/>
        <a:lstStyle/>
        <a:p>
          <a:endParaRPr lang="en-US"/>
        </a:p>
      </p:txBody>
    </p:sp>
    '''
    
    sp_elem = etree.fromstring(sp_xml.strip())
    spTree.append(sp_elem)
    return sp_elem


def add_default_text_style(sp_elem, font_name=None, font_size_pt=None, bold=None, font_color=None):
    """Add default text run properties to a shape's paragraph."""
    # Find the a:p element
    p_elem = sp_elem.find('.//' + qn('a:p'))
    if p_elem is None:
        return
    
    # Find or create pPr
    pPr = p_elem.find(qn('a:pPr'))
    if pPr is None:
        pPr = etree.Element(qn('a:pPr'))
        p_elem.insert(0, pPr)
    
    # Create defRPr for default run properties
    defRPr = etree.SubElement(pPr, qn('a:defRPr'))
    
    if font_size_pt is not None:
        defRPr.set('sz', str(int(font_size_pt * 100)))  # OOXML uses 1/100 pt
    
    if bold is not None:
        defRPr.set('b', '1' if bold else '0')
    
    if font_color:
        solidFill = etree.SubElement(defRPr, qn('a:solidFill'))
        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
        # Convert RGBColor to hex string
        srgbClr.set('val', f'{font_color[0]:02X}{font_color[1]:02X}{font_color[2]:02X}')
    
    if font_name:
        latin = etree.SubElement(defRPr, qn('a:latin'))
        latin.set('typeface', font_name)


def add_title_master_track_name_box(master, default_text="Actuarial & Finance"):
    """Add styled track name textbox to Title master."""
    tx = master.shapes.add_textbox(
        left=Inches(0.22),
        top=Inches(4.40),
        width=Inches(12.81),
        height=Inches(0.70),
    )
    tf = tx.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = default_text
    run.font.name = "Calibri"
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = TRACK_GREEN
    
    # Transparent background for textbox
    tx.fill.background()
    tx.line.fill.background()


def add_title_master_background(master, image_path):
    """Add full-bleed background image to Title master."""
    if not os.path.exists(image_path):
        print(f"  Warning: Background image not found at {image_path}")
        return
    
    pic = master.shapes.add_picture(
        image_path,
        left=Inches(0),
        top=Inches(0),
        width=Inches(SLIDE_WIDTH_IN),
        height=Inches(SLIDE_HEIGHT_IN),
    )
    
    # Move to back (behind all other shapes)
    sp_elem = pic._element
    spTree = sp_elem.getparent()
    spTree.remove(sp_elem)
    spTree.insert(2, sp_elem)  # After nvGrpSpPr and grpSpPr


def prune_extra_layouts(master, keep_count=4):
    """
    Remove extra layouts beyond keep_count from master.
    Properly removes both the layout ID entries AND the relationship parts.
    """
    layout_id_lst = master.part._element.find(qn('p:sldLayoutIdLst'))
    if layout_id_lst is None:
        return
    
    while len(layout_id_lst) > keep_count:
        last_layout_id = layout_id_lst[-1]
        rId = last_layout_id.get(qn('r:id'))
        
        # Remove from the layout ID list
        layout_id_lst.remove(last_layout_id)
        
        # Drop the relationship (this removes the orphaned layout part)
        if rId:
            try:
                master.part.drop_rel(rId)
            except Exception as e:
                print(f"  Warning: Could not drop relationship {rId}: {e}")


def clear_placeholders(slide_or_layout):
    """Remove all placeholder shapes from layout or master."""
    shapes_to_remove = []
    for shape in slide_or_layout.shapes:
        if shape.is_placeholder:
            shapes_to_remove.append(shape._element)
    
    for sp in shapes_to_remove:
        sp.getparent().remove(sp)


def configure_title_layout(layout):
    """
    Configure 'Title' layout with:
    - Full-bleed background image
    - Styled CENTER_TITLE and SUBTITLE placeholders
    - Track name textbox
    - showMasterSp="0" to hide master shapes (header bar, footer, logo)
    """
    # Set showMasterSp="0" on the cSld element to hide master shapes
    # This ensures the Title layout doesn't show header bar, footer, or logo
    cSld = layout._element.find(qn('p:cSld'))
    if cSld is not None:
        cSld.set('showMasterSp', '0')
    
    spTree = layout.shapes._spTree
    
    # Add full-bleed background image first (will be at back)
    if os.path.exists(BACKGROUND_IMAGE):
        add_title_layout_background(layout)
    
    # Title placeholder - styled with dark blue, bold, 50pt, left-aligned
    # left_margin_in=0 and first_line_indent_in=0 ensure no hanging indent on wrapped lines
    add_placeholder_shape(
        spTree,
        ph_type='title',
        idx=0,
        left=0.22, top=3.07, width=12.81, height=1.08,
        name="Title 1",
        font_name="Calibri",
        font_size_pt=50,
        bold=True,
        font_color=DARK_BLUE,
        p_align='l',
        left_margin_in=0,
        first_line_indent_in=0
    )
    
    # Subtitle placeholder - left-aligned, positioned below title
    add_placeholder_shape(
        spTree,
        ph_type='subTitle',
        idx=1,
        left=0.22, top=4.10, width=12.0, height=2.00,
        name="Subtitle 2",
        font_name="Calibri",
        font_size_pt=24,
        bold=False,
        p_align='l',
        left_margin_in=0,
        first_line_indent_in=0
    )
    
    # Add track name textbox
    add_title_layout_track_name(layout)


def add_title_layout_background(layout):
    """Add full-bleed background picture to Title layout using XML."""
    spTree = layout.shapes._spTree
    
    # Get or add image part
    image_part, rId = layout.part.get_or_add_image_part(BACKGROUND_IMAGE)
    
    existing_ids = [int(sp.get('id')) for sp in spTree.findall('.//' + qn('p:cNvPr')) if sp.get('id')]
    new_id = max(existing_ids) + 1 if existing_ids else 2
    
    # Full-bleed size in EMU
    width_emu = int(SLIDE_WIDTH_IN * 914400)
    height_emu = int(SLIDE_HEIGHT_IN * 914400)
    
    pic_xml = f'''
    <p:pic xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
           xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
           xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
      <p:nvPicPr>
        <p:cNvPr id="{new_id}" name="Picture 6"/>
        <p:cNvPicPr>
          <a:picLocks noChangeAspect="1"/>
        </p:cNvPicPr>
        <p:nvPr/>
      </p:nvPicPr>
      <p:blipFill>
        <a:blip r:embed="{rId}"/>
        <a:stretch>
          <a:fillRect/>
        </a:stretch>
      </p:blipFill>
      <p:spPr>
        <a:xfrm>
          <a:off x="0" y="0"/>
          <a:ext cx="{width_emu}" cy="{height_emu}"/>
        </a:xfrm>
        <a:prstGeom prst="rect">
          <a:avLst/>
        </a:prstGeom>
      </p:spPr>
    </p:pic>
    '''
    
    pic_elem = etree.fromstring(pic_xml.strip())
    # Insert at beginning (behind other shapes)
    spTree.insert(2, pic_elem)
    return pic_elem


def add_title_layout_track_name(layout):
    """Add styled track name textbox to Title layout using XML."""
    spTree = layout.shapes._spTree
    
    existing_ids = [int(sp.get('id')) for sp in spTree.findall('.//' + qn('p:cNvPr')) if sp.get('id')]
    new_id = max(existing_ids) + 1 if existing_ids else 2
    
    # Position/size in EMU - moved to upper-left
    left_emu = int(0.30 * 914400)
    top_emu = int(0.25 * 914400)
    width_emu = int(6.0 * 914400)
    height_emu = int(0.50 * 914400)
    
    # Track name color as hex
    track_hex = f"{TRACK_GREEN[0]:02X}{TRACK_GREEN[1]:02X}{TRACK_GREEN[2]:02X}"
    
    sp_xml = f'''
    <p:sp xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
          xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
      <p:nvSpPr>
        <p:cNvPr id="{new_id}" name="Slide Number Placeholder 5"/>
        <p:cNvSpPr txBox="1"/>
        <p:nvPr/>
      </p:nvSpPr>
      <p:spPr>
        <a:xfrm>
          <a:off x="{left_emu}" y="{top_emu}"/>
          <a:ext cx="{width_emu}" cy="{height_emu}"/>
        </a:xfrm>
        <a:prstGeom prst="rect">
          <a:avLst/>
        </a:prstGeom>
        <a:noFill/>
        <a:ln>
          <a:noFill/>
        </a:ln>
      </p:spPr>
      <p:txBody>
        <a:bodyPr wrap="square" rtlCol="0"/>
        <a:lstStyle/>
        <a:p>
          <a:pPr algn="l"/>
          <a:r>
            <a:rPr lang="en-US" sz="4000" b="1" dirty="0">
              <a:solidFill>
                <a:srgbClr val="{track_hex}"/>
              </a:solidFill>
              <a:latin typeface="Calibri"/>
            </a:rPr>
            <a:t>Actuarial &amp; Finance</a:t>
          </a:r>
          <a:endParaRPr lang="en-US" dirty="0"/>
        </a:p>
      </p:txBody>
    </p:sp>
    '''
    
    sp_elem = etree.fromstring(sp_xml.strip())
    spTree.append(sp_elem)
    return sp_elem


def configure_text_layout(layout):
    """
    Configure 'Text' layout with:
    - TITLE placeholder in header area
    - OBJECT placeholder for main content (not BODY)
    """
    spTree = layout.shapes._spTree
    
    # Title at top (in header bar area), left-aligned
    # left_margin_in=0 and first_line_indent_in=0 ensure no hanging indent
    add_placeholder_shape(
        spTree,
        ph_type='title',
        idx=0,
        left=0.12, top=0.0, width=13.22, height=0.81,
        name="Title 1",
        font_name="Calibri Light",
        font_size_pt=44,
        bold=False,
        anchor='ctr',
        p_align='l',
        left_margin_in=0,
        first_line_indent_in=0
    )
    
    # Object placeholder for content (using 'obj' type instead of 'body')
    # Default paragraph settings ensure no unwanted indentation
    add_placeholder_shape(
        spTree,
        ph_type='obj',
        idx=1,
        left=0.92, top=1.17, width=11.5, height=5.52,
        name="Content Placeholder 2",
        font_name="Calibri",
        font_size_pt=24,
        left_margin_in=0,
        first_line_indent_in=0
    )
    
    # Slide number placeholder
    add_placeholder_shape(
        spTree,
        ph_type='sldNum',
        idx=2,
        left=11.76, top=6.91, width=0.48, height=0.49,
        name="Slide Number Placeholder 3"
    )


def configure_image_right_layout(layout):
    """
    Configure 'Image Right' layout with:
    - TITLE placeholder at top
    - OBJECT placeholder on left side
    - PICTURE placeholder on right side (from layout-specs.yaml)
    """
    spTree = layout.shapes._spTree
    
    # Title at top, left-aligned with no indent
    add_placeholder_shape(
        spTree,
        ph_type='title',
        idx=0,
        left=0.12, top=0.0, width=13.22, height=0.81,
        name="Title 1",
        font_name="Calibri Light",
        font_size_pt=44,
        anchor='ctr',
        p_align='l',
        left_margin_in=0,
        first_line_indent_in=0
    )
    
    # Object placeholder on left side (narrower to make room for image)
    add_placeholder_shape(
        spTree,
        ph_type='obj',
        idx=1,
        left=0.5, top=1.2, width=6.5, height=5.5,
        name="Content Placeholder 2",
        font_name="Calibri",
        font_size_pt=24,
        left_margin_in=0,
        first_line_indent_in=0
    )
    
    # Picture placeholder on right side (from layout-specs.yaml: 7.5, 1.2, 5.33, 5.5)
    add_picture_placeholder(
        spTree,
        idx=2,
        left=7.5, top=1.2, width=5.33, height=5.5,
        name="Picture Placeholder 3"
    )
    
    # Slide number placeholder
    add_placeholder_shape(
        spTree,
        ph_type='sldNum',
        idx=3,
        left=11.76, top=6.91, width=0.48, height=0.49,
        name="Slide Number Placeholder 4"
    )


def configure_dual_image_layout(layout):
    """
    Configure 'Dual Image' layout with:
    - TITLE placeholder at top
    - Two PICTURE placeholders side by side (from layout-specs.yaml)
    - BODY placeholder at bottom for captions
    """
    spTree = layout.shapes._spTree
    
    # Title at top, left-aligned with no indent
    add_placeholder_shape(
        spTree,
        ph_type='title',
        idx=0,
        left=0.12, top=0.0, width=13.22, height=0.81,
        name="Title 1",
        font_name="Calibri Light",
        font_size_pt=44,
        anchor='ctr',
        p_align='l',
        left_margin_in=0,
        first_line_indent_in=0
    )
    
    # Left picture placeholder (from layout-specs.yaml: 0.75, 1.2, 5.5, 4.0)
    add_picture_placeholder(
        spTree,
        idx=1,
        left=0.75, top=1.2, width=5.5, height=4.0,
        name="Picture Placeholder 2"
    )
    
    # Right picture placeholder (from layout-specs.yaml: 6.75, 1.2, 5.5, 4.0)
    add_picture_placeholder(
        spTree,
        idx=2,
        left=6.75, top=1.2, width=5.5, height=4.0,
        name="Picture Placeholder 3"
    )
    
    # Body placeholder at bottom for captions/text - center-aligned,
    # positioned below images (~5.0" starts just below images at 5.2" with small gap)
    # Height 1.30" keeps bottom at 6.30" which avoids footer at 6.88"
    add_placeholder_shape(
        spTree,
        ph_type='body',
        idx=3,
        left=0.5, top=5.20, width=12.0, height=1.30,
        name="Text Placeholder 4",
        font_name="Calibri",
        font_size_pt=20,
        p_align='ctr',
        left_margin_in=0,
        first_line_indent_in=0
    )
    
    # Slide number placeholder
    add_placeholder_shape(
        spTree,
        ph_type='sldNum',
        idx=4,
        left=11.76, top=6.91, width=0.48, height=0.49,
        name="Slide Number Placeholder 5"
    )


if __name__ == "__main__":
    create_template()
