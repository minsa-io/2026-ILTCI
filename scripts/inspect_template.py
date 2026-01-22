#!/usr/bin/env python3
"""Inspect PPTX template structure: slide masters, layouts, placeholders."""

from pptx import Presentation
from pptx.util import Inches, Emu
import json
import sys

def emu_to_inches(emu):
    """Convert EMUs to inches for readability."""
    return round(emu / 914400, 2) if emu else 0

def inspect_template(template_path: str):
    """Dump template structure to stdout and JSON."""
    prs = Presentation(template_path)
    
    result = {
        "template_path": template_path,
        "slide_dimensions": {
            "width_inches": emu_to_inches(prs.slide_width),
            "height_inches": emu_to_inches(prs.slide_height),
            "width_emu": prs.slide_width,
            "height_emu": prs.slide_height,
        },
        "slide_masters": []
    }
    
    print(f"=== Template: {template_path} ===")
    print(f"Slide dimensions: {result['slide_dimensions']['width_inches']}\" x {result['slide_dimensions']['height_inches']}\"")
    print(f"  (EMU: {prs.slide_width} x {prs.slide_height})")
    print()
    
    for master_idx, master in enumerate(prs.slide_masters):
        master_data = {
            "index": master_idx,
            "name": master.name if hasattr(master, 'name') else f"Master {master_idx}",
            "layouts": []
        }
        
        print(f"--- Slide Master {master_idx}: {master_data['name']} ---")
        print(f"  Layouts count: {len(master.slide_layouts)}")
        print()
        
        for layout_idx, layout in enumerate(master.slide_layouts):
            layout_data = {
                "index": layout_idx,
                "name": layout.name,
                "placeholders": []
            }
            
            print(f"  Layout {layout_idx}: \"{layout.name}\"")
            print(f"    Placeholders ({len(layout.placeholders)}):")
            
            for ph in layout.placeholders:
                ph_data = {
                    "idx": ph.placeholder_format.idx,
                    "type": str(ph.placeholder_format.type),
                    "name": ph.name,
                    "geometry": {
                        "left_inches": emu_to_inches(ph.left),
                        "top_inches": emu_to_inches(ph.top),
                        "width_inches": emu_to_inches(ph.width),
                        "height_inches": emu_to_inches(ph.height),
                    }
                }
                layout_data["placeholders"].append(ph_data)
                
                print(f"      [{ph.placeholder_format.idx}] {ph.placeholder_format.type}: \"{ph.name}\"")
                print(f"          pos: ({ph_data['geometry']['left_inches']}\", {ph_data['geometry']['top_inches']}\") "
                      f"size: {ph_data['geometry']['width_inches']}\" x {ph_data['geometry']['height_inches']}\"")
            
            master_data["layouts"].append(layout_data)
            print()
        
        result["slide_masters"].append(master_data)
    
    return result

if __name__ == "__main__":
    template_path = sys.argv[1] if len(sys.argv) > 1 else "templates/template-alt.pptx"
    
    result = inspect_template(template_path)
    
    # Write JSON output for easier processing
    output_path = "plans/template-inspection.json"
    with open(output_path, "w") as f:
        json.dump(result, f, indent=2)
    print(f"\nJSON output written to: {output_path}")
