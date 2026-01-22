#!/usr/bin/env python3
"""
Add new slide layouts to template-alt.pptx for image-aware slides.

Creates 3 new layouts:
1. image-side: Text on left (60%), image placeholder on right (40%)
2. content-bg: Full background image + semi-transparent overlay + text
3. title-bg: Full background image + centered title overlay

python-pptx has limited layout creation capability, so we use a hybrid approach:
- Add picture shapes to existing layouts where possible
- Create placeholder-like behavior programmatically
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
import copy
import sys

# Slide dimensions (widescreen 16:9)
SLIDE_WIDTH_IN = 13.33
SLIDE_HEIGHT_IN = 7.5


def clone_layout(master, source_layout_idx, new_name):
    """
    Clone a slide layout by duplicating its XML.
    Returns the new layout.
    """
    source_layout = master.slide_layouts[source_layout_idx]
    
    # Get the slide layout part
    sldLayout_part = source_layout.part
    
    # We need to access internal XML to clone properly
    # This is a workaround since python-pptx doesn't support layout creation directly
    
    # For now, let's work with the existing layouts and add shapes to slides dynamically
    # Return None to indicate we'll use dynamic approach
    return None


def add_image_side_layout(prs):
    """
    Create image-side layout by modifying the slide layout XML.
    
    This adds a picture placeholder to the right side of a Title and Content layout.
    """
    # Find the "Title and Content" layout in master 1
    master = prs.slide_masters[1]
    base_layout = master.slide_layouts[0]  # "Title and Content"
    
    # python-pptx doesn't support adding placeholders to layouts directly,
    # so we'll need to manipulate XML or use a different approach
    
    # Alternative: Create these as shapes programmatically when building slides
    # For now, let's document the intended placeholder positions
    
    layout_spec = {
        "name": "image-side",
        "placeholders": {
            "title": {"left": 0.5, "top": 0.2, "width": 12.33, "height": 0.8},
            "body": {"left": 0.5, "top": 1.2, "width": 6.5, "height": 5.5},
            "picture": {"left": 7.5, "top": 1.2, "width": 5.33, "height": 5.5},
        }
    }
    return layout_spec


def add_content_bg_layout(prs):
    """
    Create content-bg layout specification.
    
    Background: Full-bleed image
    Foreground: Semi-transparent white overlay with title and body text
    """
    layout_spec = {
        "name": "content-bg",
        "placeholders": {
            "background": {"left": 0, "top": 0, "width": SLIDE_WIDTH_IN, "height": SLIDE_HEIGHT_IN},
            "overlay": {"left": 0.5, "top": 0.5, "width": 8.0, "height": 6.5, 
                       "fill": "white", "transparency": 0.3},
            "title": {"left": 0.75, "top": 0.75, "width": 7.5, "height": 0.8},
            "body": {"left": 0.75, "top": 1.75, "width": 7.5, "height": 5.0},
        }
    }
    return layout_spec


def add_title_bg_layout(prs):
    """
    Create title-bg layout specification.
    
    Background: Full-bleed image
    Foreground: Semi-transparent overlay strip at bottom with large title
    """
    layout_spec = {
        "name": "title-bg",
        "placeholders": {
            "background": {"left": 0, "top": 0, "width": SLIDE_WIDTH_IN, "height": SLIDE_HEIGHT_IN},
            "overlay": {"left": 0, "top": 5.0, "width": SLIDE_WIDTH_IN, "height": 2.5,
                       "fill": "black", "transparency": 0.5},
            "title": {"left": 0.5, "top": 5.25, "width": 12.33, "height": 1.5},
            "subtitle": {"left": 0.5, "top": 6.75, "width": 12.33, "height": 0.5},
        }
    }
    return layout_spec


def add_layout_via_xml(prs, layout_name, base_layout_idx=0, master_idx=1):
    """
    Add a new slide layout by cloning and modifying XML directly.
    
    This is more complex but allows us to create actual slide layouts
    that will appear in the template.
    """
    from pptx.parts.slidelayout import SlideLayoutPart
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    
    master = prs.slide_masters[master_idx]
    base_layout = master.slide_layouts[base_layout_idx]
    
    # Get the package for adding new parts
    package = prs.part.package
    
    # Clone the layout XML
    base_xml = base_layout._element
    new_xml = copy.deepcopy(base_xml)
    
    # Update the layout name in XML
    cSld = new_xml.find(qn('p:cSld'))
    if cSld is not None:
        cSld.set('name', layout_name)
    
    return new_xml  # We'd need more work to actually add this to the presentation


def create_layout_specs():
    """
    Create all layout specifications for dynamic slide building.
    
    Since python-pptx has limited layout creation ability, we'll use
    these specs to build slides dynamically with the correct shapes
    and positions.
    """
    specs = {
        "image-side": {
            "description": "Text on left (60%), image on right (40%)",
            "elements": {
                "title": {
                    "type": "placeholder",
                    "left": Inches(0.5),
                    "top": Inches(0.2),
                    "width": Inches(12.33),
                    "height": Inches(0.8)
                },
                "body": {
                    "type": "placeholder",
                    "left": Inches(0.5),
                    "top": Inches(1.2),
                    "width": Inches(6.5),
                    "height": Inches(5.5)
                },
                "picture": {
                    "type": "picture_area",
                    "left": Inches(7.5),
                    "top": Inches(1.2),
                    "width": Inches(5.33),
                    "height": Inches(5.5)
                }
            }
        },
        "content-bg": {
            "description": "Full background image with semi-transparent content overlay",
            "elements": {
                "background": {
                    "type": "background_image",
                    "left": Inches(0),
                    "top": Inches(0),
                    "width": Inches(SLIDE_WIDTH_IN),
                    "height": Inches(SLIDE_HEIGHT_IN)
                },
                "overlay": {
                    "type": "rectangle",
                    "left": Inches(0.5),
                    "top": Inches(0.5),
                    "width": Inches(8.0),
                    "height": Inches(6.5),
                    "fill_color": (255, 255, 255),
                    "transparency": 0.25
                },
                "title": {
                    "type": "textbox",
                    "left": Inches(0.75),
                    "top": Inches(0.75),
                    "width": Inches(7.5),
                    "height": Inches(0.8)
                },
                "body": {
                    "type": "textbox",
                    "left": Inches(0.75),
                    "top": Inches(1.75),
                    "width": Inches(7.5),
                    "height": Inches(5.0)
                }
            }
        },
        "title-bg": {
            "description": "Full background image with title overlay at bottom",
            "elements": {
                "background": {
                    "type": "background_image",
                    "left": Inches(0),
                    "top": Inches(0),
                    "width": Inches(SLIDE_WIDTH_IN),
                    "height": Inches(SLIDE_HEIGHT_IN)
                },
                "overlay": {
                    "type": "rectangle",
                    "left": Inches(0),
                    "top": Inches(5.0),
                    "width": Inches(SLIDE_WIDTH_IN),
                    "height": Inches(2.5),
                    "fill_color": (0, 0, 0),
                    "transparency": 0.5
                },
                "title": {
                    "type": "textbox",
                    "left": Inches(0.5),
                    "top": Inches(5.25),
                    "width": Inches(12.33),
                    "height": Inches(1.5),
                    "font_size": Pt(44),
                    "font_color": (255, 255, 255)
                },
                "subtitle": {
                    "type": "textbox",
                    "left": Inches(0.5),
                    "top": Inches(6.75),
                    "width": Inches(12.33),
                    "height": Inches(0.5),
                    "font_size": Pt(24),
                    "font_color": (255, 255, 255)
                }
            }
        }
    }
    return specs


def save_layout_specs_to_config():
    """Save layout specifications to a YAML config file."""
    import yaml
    
    # Convert Emu values to serializable format
    specs = create_layout_specs()
    
    # Convert Inches objects to float values
    def convert_emu(obj):
        if isinstance(obj, (Emu, int)) and hasattr(obj, '__int__'):
            return int(obj) / 914400  # Convert EMU to inches
        elif isinstance(obj, dict):
            return {k: convert_emu(v) for k, v in obj.items()}
        elif isinstance(obj, (list, tuple)):
            return [convert_emu(item) for item in obj]
        return obj
    
    serializable_specs = {}
    for name, spec in specs.items():
        serializable_specs[name] = {
            "description": spec["description"],
            "elements": {}
        }
        for elem_name, elem in spec["elements"].items():
            serializable_specs[name]["elements"][elem_name] = {
                k: (int(v)/914400 if isinstance(v, Emu) else v)
                for k, v in elem.items()
            }
    
    with open("assets/layout-specs.yaml", "w") as f:
        yaml.dump(serializable_specs, f, default_flow_style=False)
    
    print("Layout specs saved to assets/layout-specs.yaml")
    return serializable_specs


if __name__ == "__main__":
    # Print layout specifications
    import json
    
    print("=== Layout Specifications ===\n")
    
    specs = create_layout_specs()
    for name, spec in specs.items():
        print(f"Layout: {name}")
        print(f"  Description: {spec['description']}")
        print(f"  Elements:")
        for elem_name, elem in spec["elements"].items():
            left_in = int(elem["left"]) / 914400
            top_in = int(elem["top"]) / 914400
            width_in = int(elem["width"]) / 914400
            height_in = int(elem["height"]) / 914400
            print(f"    {elem_name}: ({left_in:.2f}\", {top_in:.2f}\") {width_in:.2f}\" x {height_in:.2f}\"")
        print()
    
    # Save to config file
    save_layout_specs_to_config()
    
    print("\nThese specifications will be used dynamically when building slides.")
    print("The generator will apply these layouts based on slide directives.")
