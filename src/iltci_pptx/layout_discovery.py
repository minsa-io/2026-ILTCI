"""Layout discovery and registry for PowerPoint templates.

This module discovers slide layouts from a PowerPoint template and provides
a registry mapping layout names to their indices. This enables template-agnostic
slide generation by using layout names instead of hardcoded indices.
"""

import logging
from pathlib import Path
from typing import Union

from pptx import Presentation

logger = logging.getLogger(__name__)

# Type alias: maps layout name to its index in prs.slide_layouts
LayoutRegistry = dict[str, int]


def _has_placeholders(layout) -> bool:
    """Check if a slide layout has any placeholders.
    
    Layouts with placeholders are considered "master layouts" suitable
    for content placement (title, body text, images, etc.).
    
    Args:
        layout: A SlideLayout object from python-pptx.
        
    Returns:
        True if the layout has at least one placeholder.
    """
    # placeholders is a collection; check if it has any items
    try:
        return len(layout.placeholders) > 0
    except (AttributeError, TypeError):
        return False


def load_layout_registry(template_path: Union[str, Path]) -> LayoutRegistry:
    """Load and build a layout registry from a PowerPoint template.
    
    Discovers all slide layouts in the template and returns a dictionary
    mapping layout names to their indices in `prs.slide_layouts`.
    
    The registry only includes layouts that have placeholders, as these
    are the layouts suitable for content placement. Layouts without
    placeholders (e.g., blank backgrounds) are excluded with a debug log.
    
    If duplicate layout names are found (which can happen with multiple
    slide masters), a warning is logged and only the first layout with
    that name is kept.
    
    Args:
        template_path: Path to the PowerPoint template file (.pptx).
        
    Returns:
        Dictionary mapping layout name to layout index.
        
    Raises:
        FileNotFoundError: If the template file does not exist.
        
    Example:
        >>> registry = load_layout_registry("templates/template.pptx")
        >>> print(registry)
        {'Title Slide': 0, 'Title and Content': 1, 'Two Content': 2, ...}
        >>> # Use with python-pptx:
        >>> prs = Presentation("templates/template.pptx")
        >>> layout = prs.slide_layouts[registry['Title and Content']]
    """
    path = Path(template_path)
    if not path.exists():
        raise FileNotFoundError(f"Template file not found: {template_path}")
    
    prs = Presentation(str(path))
    registry: LayoutRegistry = {}
    
    logger.info(f"Discovering layouts from template: {path}")
    
    # Iterate through all slide layouts in the presentation
    # Note: prs.slide_layouts contains layouts from all masters, flattened
    for idx, layout in enumerate(prs.slide_layouts):
        layout_name = layout.name
        
        # Check if layout has placeholders (suitable for content)
        if not _has_placeholders(layout):
            logger.debug(
                f"Skipping layout '{layout_name}' (index {idx}): no placeholders"
            )
            continue
        
        # Check for duplicate names
        if layout_name in registry:
            existing_idx = registry[layout_name]
            logger.warning(
                f"Duplicate layout name '{layout_name}' found at index {idx}; "
                f"keeping first occurrence at index {existing_idx}"
            )
            continue
        
        registry[layout_name] = idx
        logger.debug(
            f"Registered layout '{layout_name}' at index {idx} "
            f"({len(layout.placeholders)} placeholders)"
        )
    
    logger.info(
        f"Layout registry built: {len(registry)} layouts discovered"
    )
    
    if logger.isEnabledFor(logging.DEBUG):
        logger.debug(f"Available layouts: {list(registry.keys())}")
    
    return registry


def get_available_layout_names(registry: LayoutRegistry) -> list[str]:
    """Get a sorted list of available layout names from a registry.
    
    Useful for displaying available options in error messages or CLI help.
    
    Args:
        registry: A layout registry from load_layout_registry().
        
    Returns:
        Sorted list of layout names.
    """
    return sorted(registry.keys())


def validate_layout_name(
    layout_name: str,
    registry: LayoutRegistry,
    raise_on_missing: bool = True
) -> bool:
    """Validate that a layout name exists in the registry.
    
    Args:
        layout_name: The layout name to validate.
        registry: A layout registry from load_layout_registry().
        raise_on_missing: If True, raise ValueError when layout not found.
            If False, return False instead.
            
    Returns:
        True if layout exists in registry.
        
    Raises:
        ValueError: If layout not found and raise_on_missing is True.
            Error message includes list of available layouts.
    """
    if layout_name in registry:
        return True
    
    if raise_on_missing:
        available = get_available_layout_names(registry)
        raise ValueError(
            f"Unknown layout '{layout_name}'. "
            f"Available layouts: {', '.join(available)}"
        )
    
    return False
