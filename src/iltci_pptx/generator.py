"""Main presentation generation orchestration.

This module provides the primary interface for generating PowerPoint presentations
from markdown content using a template-agnostic, registry-driven pipeline.

Pipeline flow:
    1. Load config and layout registry from template
    2. Parse markdown into SlideData objects (with layout validation)
    3. For each slide: build_slide → populate_slide (includes images)
    4. Save the presentation
"""

import logging
from pathlib import Path
from pptx import Presentation
from typing import Optional

from .config import Config
from .layout_discovery import load_layout_registry, LayoutRegistry, get_available_layout_names
from .markdown_parser import parse_markdown_file, parse_markdown_slides, SlideData
from .slide_builders import build_slide, populate_slide

logger = logging.getLogger(__name__)


class PresentationGenerator:
    """Orchestrates the creation of PowerPoint presentations from markdown.
    
    Uses a registry-driven pipeline for template-agnostic slide generation:
    - Layout discovery builds a name→index registry from the template
    - Markdown parsing validates layouts against the registry
    - Unified build/populate functions handle all slide types generically
    """
    
    def __init__(self, config: Config):
        """Initialize the generator with configuration.
        
        Args:
            config: Configuration object with paths and settings.
        """
        self.config = config
        self._registry: Optional[LayoutRegistry] = None
    
    @property
    def registry(self) -> LayoutRegistry:
        """Lazy-load the layout registry from the template.
        
        Returns:
            LayoutRegistry mapping layout names to indices.
        """
        if self._registry is None:
            self._registry = load_layout_registry(self.config.template_path)
        return self._registry
    
    def generate(self, template_override: Optional[Path] = None) -> None:
        """Generate the PowerPoint presentation from markdown content.
        
        Uses the new generic pipeline:
        1. Load layout registry from template
        2. Parse markdown with layout validation
        3. Build and populate slides using registry-driven functions
        
        Args:
            template_override: Optional path to override the template from config.
            style_overrides: Optional dict of style overrides to apply to slides.
        """
        # Validate paths exist
        self.config.validate_paths()
        
        # Determine template path (may be overridden or set via frontmatter later)
        template_path = template_override or self.config.template_path
        
        # Load layout registry from template (for validation and building)
        logger.info(f"Discovering layouts from template: {template_path}")
        registry = load_layout_registry(template_path)
        
        available_layouts = get_available_layout_names(registry)
        logger.info(f"Available layouts ({len(available_layouts)}): {', '.join(available_layouts)}")
        
        # Parse markdown content with layout validation
        content_path = self.config.content_path
        logger.info(f"Parsing markdown file: {content_path}")
        
        try:
            # Use new parse_markdown_file with registry validation
            doc_frontmatter, slide_data_list = parse_markdown_file(
                content_path,
                registry,
                self.config,
                strict=False,  # Warn on invalid layouts instead of failing
            )
        except ValueError as e:
            # If strict parsing failed, fall back to legacy parser
            logger.warning(f"New parser encountered issues: {e}")
            doc_frontmatter = {}
            slide_data_list = []
        
        logger.info(f"Document frontmatter keys: {list(doc_frontmatter.keys())}")
        
        # Fall back to legacy parser if new parser found no slides
        # (likely means markdown doesn't have per-slide frontmatter)
        if not slide_data_list:
            logger.info("No slides parsed with new format - attempting legacy fallback...")
            slide_data_list = self._parse_legacy_fallback(content_path, registry)
        
        logger.info(f"Parsed {len(slide_data_list)} slides")
        
        # Check for template override in document frontmatter
        # NOTE: We don't reload the registry because slides are already validated
        # against the current registry. The frontmatter template should match the
        # config template for consistent behavior.
        if 'template' in doc_frontmatter and not template_override:
            new_template = self.config.project_root / doc_frontmatter['template']
            if new_template != template_path:
                logger.info(f"Template in frontmatter: {new_template}")
                # Verify the new template has compatible layouts
                new_registry = load_layout_registry(new_template)
                new_available = get_available_layout_names(new_registry)
                if new_available:
                    logger.info(f"Using frontmatter template with {len(new_available)} layouts")
                    template_path = new_template
                    registry = new_registry
                else:
                    logger.warning(
                        f"Frontmatter template '{new_template}' has no usable layouts. "
                        f"Keeping original template: {template_path}"
                    )
        
        # Load presentation template
        logger.info(f"Loading presentation template: {template_path}")
        prs = Presentation(str(template_path))
        
        logger.info(f"Template has {len(prs.slide_masters)} slide master(s)")
        logger.info(f"Template has {len(prs.slide_layouts)} total layouts")
        logger.info(f"Template has {len(prs.slides)} existing slides")
        
        # Remove existing content slides
        logger.info("Removing existing slides from template...")
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
        
        logger.info(f"Building {len(slide_data_list)} new slides...")
        
        # Build and populate slides using generic pipeline
        for idx, data in enumerate(slide_data_list):
            logger.info(f"\n=== Slide {idx + 1}: {data.layout_name} ===")
            if data.title:
                logger.info(f"  Title: {data.title[:60]}{'...' if len(data.title) > 60 else ''}")
            if data.images:
                logger.info(f"  Images: {len(data.images)}")
            
            try:
                # Build slide using layout from registry
                slide = build_slide(data.layout_name, prs, registry)
                
                # Populate slide with content (title, body, images)
                populate_slide(slide, data, self.config, registry)
                
                logger.debug(f"  Successfully built slide {idx + 1}")
                
            except ValueError as e:
                logger.error(f"  Failed to build slide {idx + 1}: {e}")
                # Continue with remaining slides
                continue
            except Exception as e:
                logger.error(f"  Unexpected error on slide {idx + 1}: {e}")
                raise
        
        # Save presentation
        output_path = self.config.output_path
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        logger.info(f"\nSaving presentation to {output_path}...")
        prs.save(str(output_path))
        logger.info("✓ Presentation saved successfully!")
        logger.info(f"  Total slides created: {len(prs.slides)}")
    
    def _parse_legacy_fallback(
        self,
        content_path: Path,
        registry: LayoutRegistry
    ) -> list[SlideData]:
        """Fall back to legacy parser and convert dict format to SlideData.
        
        This provides backward compatibility for markdown files that don't
        use the new per-slide frontmatter format.
        
        Args:
            content_path: Path to markdown file.
            registry: Layout registry for inferring layout names.
            
        Returns:
            List of SlideData objects converted from legacy dict format.
        """
        logger.info("Using legacy parser for backward compatibility")
        
        # Use legacy parser (returns dict-based slides)
        frontmatter, legacy_slides = parse_markdown_slides(content_path, self.config)
        
        # Convert legacy dicts to SlideData objects
        slide_data_list: list[SlideData] = []
        available = get_available_layout_names(registry)
        
        for idx, slide_dict in enumerate(legacy_slides):
            # Determine layout name
            layout_name = slide_dict.get('layout')
            
            if not layout_name:
                # Infer layout from slide type
                if slide_dict.get('is_title'):
                    # Try common title slide layout names
                    for candidate in ['Title Slide', 'title-slide', 'Title', 'title']:
                        if candidate in registry:
                            layout_name = candidate
                            break
                else:
                    # Try common content layout names
                    for candidate in ['Title and Content', 'content', 'Content', 'Body']:
                        if candidate in registry:
                            layout_name = candidate
                            break
            
            # Validate layout exists
            if not layout_name or layout_name not in registry:
                logger.warning(
                    f"Slide {idx + 1}: Could not determine valid layout "
                    f"(tried: {layout_name}). Available: {', '.join(available[:5])}..."
                )
                # Use first available layout as fallback
                layout_name = available[0] if available else 'Title and Content'
            
            # Build content blocks from legacy format
            content_blocks = []
            if slide_dict.get('content'):
                content_blocks = [slide_dict['content']]
            elif slide_dict.get('subtitle'):
                content_blocks = [slide_dict['subtitle']]
            
            # Extract images from legacy format
            images = []
            bg_image = slide_dict.get('bg_image')
            if bg_image:
                images.append(bg_image)
            
            # Build SlideData
            slide_data = SlideData(
                layout_name=layout_name,
                title=slide_dict.get('title'),
                content_blocks=content_blocks,
                images=images,
                section_name=slide_dict.get('section_name', ''),
                options={
                    k: v for k, v in slide_dict.items()
                    if k not in ('is_title', 'title', 'subtitle', 'content', 
                                'layout', 'section_name', 'bg_image')
                },
            )
            
            slide_data_list.append(slide_data)
        
        logger.info(f"Converted {len(slide_data_list)} slides from legacy format")
        return slide_data_list
