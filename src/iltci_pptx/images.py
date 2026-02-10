"""Image handling for PowerPoint slides.

This module provides functions for adding images to slides, using PICTURE
placeholders discovered from the template. No external YAML configuration
is required - image placement is driven entirely by template placeholders.
"""

import re
import logging
from pathlib import Path
from pptx.util import Inches, Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE as PH_TYPE
from typing import List, Dict, Any, Optional, Tuple, TYPE_CHECKING
from pptx.enum.text import PP_ALIGN
from .config import Config
from .markdown_parser import SlideData

if TYPE_CHECKING:
    from pptx.slide import Slide
    from pptx.shapes.picture import Picture
    from .layout_discovery import LayoutRegistry

# Default image style settings for borders and rounded corners
IMAGE_STYLE_DEFAULTS = {
    'border_width': Pt(2),           # 2pt border width
    'border_color': RGBColor(68, 84, 106),  # #44546A - dark blue-gray
    'corner_radius': Inches(0.1),    # Rounded corner radius
    'border_enabled': True,          # Enable border by default
    'rounded_enabled': True,         # Enable rounded corners by default
}

# Style class mappings for per-image overrides
STYLE_CLASS_MAP = {
    'no-border': {'border_enabled': False},
    'no-rounded': {'rounded_enabled': False},
    'border-thin': {'border_width': Pt(1)},
    'border-thick': {'border_width': Pt(4)},
    'border-light': {'border_color': RGBColor(180, 198, 231)},  # #B4C6E7
    'border-dark': {'border_color': RGBColor(47, 62, 78)},      # #2F3E4E
    'rounded-sm': {'corner_radius': Inches(0.05)},
    'rounded-lg': {'corner_radius': Inches(0.2)},
}


def get_picture_placeholders(slide: 'Slide') -> List[Any]:
    """Get all PICTURE placeholders from a slide, sorted by left position.
    
    This function discovers image placement areas from the slide's layout
    by finding all placeholders with type PICTURE. The placeholders are
    returned sorted by their left position for consistent ordering.
    
    Args:
        slide: PowerPoint slide object
        
    Returns:
        List of placeholder shapes with type PICTURE, sorted left-to-right
    """
    picture_phs = []
    for shape in slide.shapes:
        if not hasattr(shape, 'placeholder_format'):
            continue
        ph_format = shape.placeholder_format
        if ph_format and ph_format.type == PH_TYPE.PICTURE:
            picture_phs.append(shape)
    
    # Sort by left position for consistent left-to-right ordering
    picture_phs.sort(key=lambda ph: ph.left)
    return picture_phs


def parse_style_classes(class_attr: str) -> Dict[str, Any]:
    """Parse CSS class attribute and return merged style settings.
    
    Args:
        class_attr: Space-separated CSS class names from HTML img tag
        
    Returns:
        Dictionary of style settings merged from defaults and class overrides
    """
    style = IMAGE_STYLE_DEFAULTS.copy()
    
    if not class_attr:
        return style
    
    classes = class_attr.split()
    for cls in classes:
        if cls in STYLE_CLASS_MAP:
            style.update(STYLE_CLASS_MAP[cls])
    
    return style


def compute_image_style(
    class_attr: Optional[str] = None,
    per_image_src_overrides: Optional[Dict[str, Any]] = None
) -> Dict[str, Any]:
    """Compute final image style by merging defaults, class overrides, and per-image overrides.
    
    Precedence (each layer overrides the previous):
    1. IMAGE_STYLE_DEFAULTS (base)
    2. STYLE_CLASS_MAP classes from class_attr
    3. per_image_src_overrides from config (highest precedence)
    
    Args:
        class_attr: Space-separated CSS class names from markdown image metadata
        per_image_src_overrides: Per-image style overrides from config.image_styles.per_image_src
        
    Returns:
        Dictionary of final style settings
    """
    # Start with defaults + class overrides
    style = parse_style_classes(class_attr or '')
    
    # Apply per-image overrides (highest precedence)
    if per_image_src_overrides:
        style.update(per_image_src_overrides)
    
    return style


def apply_image_style(picture: 'Picture', style: Dict[str, Any]) -> None:
    """Apply border and rounded corner styling to a PowerPoint Picture shape.
    
    Args:
        picture: PowerPoint Picture shape object
        style: Style dictionary with border_enabled, border_width, border_color,
               rounded_enabled, corner_radius keys
    """
    if picture is None:
        return
    
    try:
        # Apply border styling
        if style.get('border_enabled', True):
            line = picture.line
            line.width = style.get('border_width', Pt(2))
            line.color.rgb = style.get('border_color', RGBColor(68, 84, 106))
        else:
            # Remove border
            picture.line.fill.background()
        
        # Apply rounded corners using adjustments
        # PowerPoint uses shape adjustments (adj) for rounded rectangles
        # For pictures, we need to use XML manipulation for soft edges/rounded corners
        if style.get('rounded_enabled', True):
            corner_radius = style.get('corner_radius', Inches(0.1))
            _apply_rounded_corners(picture, corner_radius)
        
        logging.debug(f"Applied image style: border={style.get('border_enabled')}, rounded={style.get('rounded_enabled')}")
        
    except Exception as e:
        logging.warning(f"Could not apply image style: {e}")


def _apply_rounded_corners(picture: 'Picture', radius: int) -> None:
    """Apply rounded corners to a picture using XML manipulation.
    
    Args:
        picture: PowerPoint Picture shape
        radius: Corner radius in EMUs
    """
    try:
        from pptx.oxml.ns import qn
        from lxml import etree
        
        # Get the spPr (shape properties) element
        spPr = picture._pic.spPr
        
        # Check if prstGeom exists, if not create it
        prstGeom = spPr.find(qn('a:prstGeom'))
        if prstGeom is None:
            # Create preset geometry for rounded rectangle
            prstGeom = etree.SubElement(spPr, qn('a:prstGeom'))
        
        # Set to rounded rectangle preset
        prstGeom.set('prst', 'roundRect')
        
        # Add or update adjustment values for corner radius
        avLst = prstGeom.find(qn('a:avLst'))
        if avLst is None:
            avLst = etree.SubElement(prstGeom, qn('a:avLst'))
        
        # Clear existing adjustments
        for child in list(avLst):
            avLst.remove(child)
        
        # Calculate adjustment value (PowerPoint uses 0-50000 scale for roundRect)
        # Convert radius to percentage of shortest side
        # Assuming ~10% rounding as default
        adj_val = 10000  # 10% rounding
        if isinstance(radius, int):  # If EMUs provided
            # Rough conversion - smaller radius = smaller adj value
            adj_val = min(int(radius / 914400 * 100000), 50000)  # Cap at 50%
        
        gd = etree.SubElement(avLst, qn('a:gd'))
        gd.set('name', 'adj')
        gd.set('fmla', f'val {adj_val}')
        
        logging.debug(f"Applied rounded corners with adj={adj_val}")
        
    except ImportError:
        logging.warning("lxml not available, cannot apply rounded corners via XML")
    except Exception as e:
        logging.debug(f"Could not apply rounded corners: {e}")


def add_background_image(slide: 'Slide', img_path: Path, 
                         width: float = 13.33, height: float = 7.5) -> None:
    """Add a full-bleed background image to slide.
    
    The image is added at position (0,0) and should be moved to back.
    
    Args:
        slide: PowerPoint slide object
        img_path: Path to the image file
        width: Slide width in inches (default: 13.33 for widescreen)
        height: Slide height in inches (default: 7.5)
    """
    if not img_path.exists():
        logging.warning(f"Background image not found: {img_path}")
        return
    
    try:
        # Add picture at full slide size
        picture = slide.shapes.add_picture(
            str(img_path),
            Inches(0),
            Inches(0),
            width=Inches(width),
            height=Inches(height)
        )
        
        # Move to back by adjusting z-order
        # In python-pptx, we need to move the shape to the beginning of shapes
        sp = picture._element
        sp.getparent().insert(0, sp)
        
        logging.info(f"Added background image: {img_path}")
    except Exception as e:
        logging.error(f"Error adding background image {img_path}: {e}")


def add_overlay_rectangle(slide: 'Slide', left: float, top: float, 
                          width: float, height: float,
                          fill_color: tuple = (255, 255, 255),
                          transparency: float = 0.25) -> None:
    """Add a semi-transparent overlay rectangle to slide.
    
    Args:
        slide: PowerPoint slide object
        left, top, width, height: Position and size in inches
        fill_color: RGB tuple (r, g, b) with values 0-255
        transparency: Transparency value 0.0 (opaque) to 1.0 (fully transparent)
    """
    try:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height)
        )
        
        # Set fill color
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*fill_color)
        
        # Set transparency (python-pptx uses 0-100000 scale for transparency)
        # transparency 0.25 = 25% transparent = 75% opaque
        fill.fore_color.brightness = 0  # No brightness adjustment
        
        # Note: python-pptx doesn't directly support fill transparency easily
        # We need to use XML manipulation for proper transparency
        from pptx.oxml.ns import qn
        spPr = shape._sp.spPr
        solidFill = spPr.find(qn('a:solidFill'))
        if solidFill is not None:
            srgbClr = solidFill.find(qn('a:srgbClr'))
            if srgbClr is not None:
                # Add alpha element for transparency
                # alpha value is 0-100000 where 100000 = fully opaque
                alpha_val = int((1 - transparency) * 100000)
                alpha = srgbClr.makeelement(qn('a:alpha'), {})
                alpha.set('val', str(alpha_val))
                srgbClr.append(alpha)
        
        # Remove outline
        shape.line.fill.background()
        
        logging.info(f"Added overlay rectangle at ({left}, {top}) with {transparency*100}% transparency")
    except Exception as e:
        logging.error(f"Error adding overlay rectangle: {e}")


def add_image_to_area(slide: 'Slide', img_path: Path,
                      left: float, top: float, width: float, height: float,
                      fit_mode: str = 'contain',
                      class_attr: Optional[str] = None,
                      per_image_src_overrides: Optional[Dict[str, Any]] = None) -> Optional['Picture']:
    """Add an image to a specific area, scaling to fit.
    
    Args:
        slide: PowerPoint slide object
        img_path: Path to image file
        left, top, width, height: Target area in inches
        fit_mode: 'contain' (fit within, preserve aspect) or 'cover' (fill, may crop)
        class_attr: CSS class attribute for style overrides (e.g., 'no-border rounded-lg')
        per_image_src_overrides: Per-image style overrides from config (highest precedence)
        
    Returns:
        Picture shape or None if failed
    """
    if not img_path.exists():
        logging.warning(f"Image not found: {img_path}")
        return None
    
    # Compute final style: defaults → class overrides → per-image overrides
    image_style = compute_image_style(class_attr, per_image_src_overrides)
    
    try:
        from PIL import Image
        
        # Get original image dimensions
        with Image.open(img_path) as img:
            orig_width, orig_height = img.size
        
        # Calculate aspect ratios
        orig_ratio = orig_width / orig_height
        target_ratio = width / height
        
        if fit_mode == 'contain':
            # Scale to fit within bounds, preserving aspect ratio
            if orig_ratio > target_ratio:
                # Image is wider - constrain by width
                final_width = width
                final_height = width / orig_ratio
            else:
                # Image is taller - constrain by height
                final_height = height
                final_width = height * orig_ratio
            
            # Center within target area
            final_left = left + (width - final_width) / 2
            final_top = top + (height - final_height) / 2
        else:
            # Cover mode - fill area (may exceed bounds)
            if orig_ratio > target_ratio:
                # Image is wider - constrain by height, center horizontally
                final_height = height
                final_width = height * orig_ratio
                final_top = top
                final_left = left + (width - final_width) / 2
            else:
                # Image is taller - constrain by width, center vertically
                final_width = width
                final_height = width / orig_ratio
                final_left = left
                final_top = top + (height - final_height) / 2
        
        picture = slide.shapes.add_picture(
            str(img_path),
            Inches(final_left),
            Inches(final_top),
            width=Inches(final_width),
            height=Inches(final_height)
        )
        
        # Apply image styling (border and rounded corners)
        apply_image_style(picture, image_style)
        
        logging.info(f"Added image {img_path} to area ({left}, {top}) {width}x{height} [mode={fit_mode}]")
        return picture
        
    except ImportError:
        # PIL not available, use simple placement
        logging.warning("PIL not available, using simple image placement")
        picture = slide.shapes.add_picture(
            str(img_path),
            Inches(left),
            Inches(top),
            width=Inches(width)
        )
        
        # Apply image styling (border and rounded corners)
        apply_image_style(picture, image_style)
        return picture
        
    except Exception as e:
        logging.error(f"Error adding image to area: {e}")
        return None


def add_images_for_layout(
    slide_data: SlideData,
    slide: 'Slide',
    config: Config,
    registry: 'LayoutRegistry',
    base_path: Optional[Path] = None,
    fit_mode: str = 'contain'
) -> None:
    """Add images to slide using PICTURE placeholders from the template.
    
    Discovers PICTURE placeholders from the slide's layout and uses their
    positions to place images. No external YAML configuration required -
    the template's embedded placeholders define all image placement areas.
    
    Args:
        slide_data: SlideData containing images list and layout_name
        slide: PowerPoint slide object
        config: Configuration object for assets_dir and image_styles
        registry: LayoutRegistry (unused, kept for API compatibility)
        base_path: Base path for resolving image paths (defaults to config.assets_dir)
        fit_mode: 'contain' or 'cover' for image fitting
    """
    if not slide_data.images:
        return
    
    # Resolve base path for images
    if base_path is None:
        base_path = config.assets_dir
    
    layout_name = slide_data.layout_name
    
    # Discover PICTURE placeholders from the slide's layout
    picture_phs = get_picture_placeholders(slide)
    
    if not picture_phs:
        # No PICTURE placeholders in this layout - images cannot be placed
        logging.warning(
            f"Layout '{layout_name}' has no PICTURE placeholders in template. "
            f"{len(slide_data.images)} image(s) will not be placed."
        )
        return
    
    # Warn if more images than placeholders
    if len(slide_data.images) > len(picture_phs):
        logging.warning(
            f"Layout '{layout_name}' has {len(picture_phs)} PICTURE placeholder(s) but "
            f"{len(slide_data.images)} images were provided. "
            f"Only the first {len(picture_phs)} image(s) will be placed."
        )
    
    # Get per-image overrides from config
    per_image_src_config = config.get('image_styles.per_image_src', {}) or {}
    
    # Place images into PICTURE placeholders
    for i, img_info in enumerate(slide_data.images[:len(picture_phs)]):
        ph = picture_phs[i]
        
        img_src = img_info.get('src', '')
        if not img_src:
            logging.warning(f"Image at index {i} has no 'src' attribute, skipping")
            continue
        
        img_path = base_path / img_src
        
        if not img_path.exists():
            logging.warning(f"Image not found: {img_path}")
            continue
        
        # Look up per-image style overrides by src
        per_image_overrides = per_image_src_config.get(img_src, {})
        
        # Get placeholder position and size (convert EMU to inches)
        left = ph.left / 914400  # EMU to inches
        top = ph.top / 914400
        width = ph.width / 914400
        height = ph.height / 914400
        
        class_attr = img_info.get('class', '')
        caption = img_info.get('data-caption', '')
        
        logging.debug(
            f"Placing image into placeholder '{ph.name}' at "
            f"({left:.2f}\", {top:.2f}\") size {width:.2f}\" x {height:.2f}\""
        )
        
        if caption:
            add_image_with_caption(
                slide, img_path, left, top, width, height,
                caption=caption, fit_mode=fit_mode, class_attr=class_attr,
                per_image_src_overrides=per_image_overrides
            )
        else:
            add_image_to_area(
                slide, img_path, left, top, width, height,
                fit_mode=fit_mode, class_attr=class_attr,
                per_image_src_overrides=per_image_overrides
            )
        
        logging.info(f"Added image to PICTURE placeholder '{ph.name}'")


# Caption layout constants
CAPTION_HEIGHT = 0.6      # Caption text box height in inches
CAPTION_GAP = 0.05        # Gap between image and caption in inches

# Default caption style settings
CAPTION_STYLE = {
    'color': RGBColor(51, 51, 51),  # #333333 - off-black
    'font_size': 12,  # 12pt
    'align': PP_ALIGN.LEFT,
}


def get_caption_dimensions(image_bottom: float) -> tuple[float, float]:
    """Calculate caption position and height based on image bottom position.
    
    Args:
        image_bottom: Bottom position of the image in inches
        
    Returns:
        Tuple of (caption_top, caption_bottom) positions in inches
    """
    caption_top = image_bottom + CAPTION_GAP
    caption_bottom = caption_top + CAPTION_HEIGHT
    return caption_top, caption_bottom


def add_image_caption(slide: 'Slide', caption: str, 
                      left: float, top: float, width: float,
                      style: Optional[Dict[str, Any]] = None) -> None:
    """Add a caption text box below an image.
    
    Args:
        slide: PowerPoint slide object
        caption: Caption text to display
        left: Left position in inches (should match image left)
        top: Top position in inches (should be below image bottom)
        width: Width in inches (should match image width)
        style: Optional style overrides {color, font_size, align}
    """
    if not caption:
        return
    
    # Merge default style with any overrides
    caption_style = CAPTION_STYLE.copy()
    if style:
        if 'color' in style:
            caption_style['color'] = style['color']
        if 'font_size' in style:
            caption_style['font_size'] = style['font_size']
        if 'align' in style:
            caption_style['align'] = style['align']
    
    try:
        # Create caption text box using constant height
        caption_box = slide.shapes.add_textbox(
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(CAPTION_HEIGHT)
        )
        
        caption_frame = caption_box.text_frame
        caption_frame.word_wrap = True
        
        # Split caption text into lines (interpret \n as newlines, ; as separator)
        caption_text = caption.replace('\\n', '\n')
        # Also split on semicolon followed by space for inline multi-line captions
        caption_text = caption_text.replace('; ', '\n')
        lines = caption_text.split('\n')
        
        # Add each line as a separate paragraph
        for i, line in enumerate(lines):
            line = line.strip()
            if not line:
                continue
            
            if i == 0:
                # Use existing first paragraph
                p = caption_frame.paragraphs[0]
            else:
                # Add new paragraph for subsequent lines
                p = caption_frame.add_paragraph()
            
            p.text = line
            p.alignment = caption_style['align']
            
            # Style the text
            for run in p.runs:
                run.font.size = Pt(caption_style['font_size'])
                run.font.color.rgb = caption_style['color']
        
        # Remove any fill/background on the textbox
        caption_box.fill.background()
        
        logging.info(f"Added caption with {len(lines)} line(s): '{caption}' at ({left}, {top})")
        
    except Exception as e:
        logging.error(f"Error adding caption: {e}")


def add_image_with_caption(slide: 'Slide', img_path: Path,
                           left: float, top: float, width: float, height: float,
                           caption: Optional[str] = None,
                           fit_mode: str = 'contain',
                           caption_style: Optional[Dict[str, Any]] = None,
                           class_attr: Optional[str] = None,
                           per_image_src_overrides: Optional[Dict[str, Any]] = None) -> Tuple[Optional['Picture'], float]:
    """Add an image to a specific area with optional caption below.
    
    Args:
        slide: PowerPoint slide object
        img_path: Path to image file
        left, top, width, height: Target area in inches
        caption: Optional caption text to display below image
        fit_mode: 'contain' (fit within, preserve aspect) or 'cover' (fill, may crop)
        caption_style: Optional style overrides for caption
        class_attr: CSS class attribute for style overrides (e.g., 'no-border rounded-lg')
        per_image_src_overrides: Per-image style overrides from config (highest precedence)
        
    Returns:
        Tuple of (picture shape or None, bottom position in inches - caption bottom if 
        caption was added, otherwise image bottom). This allows callers to position
        content below the image+caption combination without overlap.
    """
    if not img_path.exists():
        logging.warning(f"Image not found: {img_path}")
        return None, top + height
    
    # Compute final style: defaults → class overrides → per-image overrides
    image_style = compute_image_style(class_attr, per_image_src_overrides)
    
    try:
        from PIL import Image
        
        # Get original image dimensions
        with Image.open(img_path) as img:
            orig_width, orig_height = img.size
        
        # Calculate aspect ratios
        orig_ratio = orig_width / orig_height
        target_ratio = width / height
        
        if fit_mode == 'contain':
            # Scale to fit within bounds, preserving aspect ratio
            if orig_ratio > target_ratio:
                # Image is wider - constrain by width
                final_width = width
                final_height = width / orig_ratio
            else:
                # Image is taller - constrain by height
                final_height = height
                final_width = height * orig_ratio
            
            # Center within target area
            final_left = left + (width - final_width) / 2
            final_top = top + (height - final_height) / 2
        else:
            # Cover mode - fill area (may exceed bounds)
            if orig_ratio > target_ratio:
                final_height = height
                final_width = height * orig_ratio
                final_top = top
                final_left = left + (width - final_width) / 2
            else:
                final_width = width
                final_height = width / orig_ratio
                final_left = left
                final_top = top + (height - final_height) / 2
        
        picture = slide.shapes.add_picture(
            str(img_path),
            Inches(final_left),
            Inches(final_top),
            width=Inches(final_width),
            height=Inches(final_height)
        )
        
        # Apply image styling (border and rounded corners)
        apply_image_style(picture, image_style)
        
        # Calculate actual bottom of image
        actual_bottom = final_top + final_height
        
        logging.info(f"Added image {img_path} to area ({left}, {top}) {width}x{height} [mode={fit_mode}]")
        
        # Add caption if provided
        if caption:
            # Use helper to calculate caption position
            caption_top, caption_bottom = get_caption_dimensions(actual_bottom)
            add_image_caption(slide, caption, final_left, caption_top, final_width, caption_style)
            return picture, caption_bottom
        
        return picture, actual_bottom
        
    except ImportError:
        # PIL not available, use simple placement
        logging.warning("PIL not available, using simple image placement")
        picture = slide.shapes.add_picture(
            str(img_path),
            Inches(left),
            Inches(top),
            width=Inches(width)
        )
        
        # Apply image styling (border and rounded corners)
        apply_image_style(picture, image_style)
        
        actual_bottom = top + height
        
        if caption:
            # Use helper to calculate caption position
            caption_top, caption_bottom = get_caption_dimensions(actual_bottom)
            add_image_caption(slide, caption, left, caption_top, width, caption_style)
            return picture, caption_bottom
        
        return picture, actual_bottom
        
    except Exception as e:
        logging.error(f"Error adding image to area: {e}")
        return None, top + height
