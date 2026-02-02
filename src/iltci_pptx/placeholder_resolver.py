"""Placeholder resolution utilities for PowerPoint slides.

This module provides utilities to resolve placeholders by name or type,
enabling generic slide population without hardcoded shape indices.

Typical usage:
    >>> from pptx import Presentation
    >>> from pptx.enum.shapes import PP_PLACEHOLDER_TYPE as PH_TYPE
    >>> from iltci_pptx.placeholder_resolver import get_placeholder, resolve_placeholders
    >>> 
    >>> prs = Presentation("templates/template.pptx")
    >>> slide = prs.slides.add_slide(prs.slide_layouts[0])
    >>> 
    >>> # Get by type
    >>> title_ph = get_placeholder(slide, ph_type=PH_TYPE.TITLE)
    >>> 
    >>> # Get by name
    >>> content_ph = get_placeholder(slide, name="Content Placeholder 2")
    >>> 
    >>> # Resolve multiple at once
    >>> specs = {"title": PH_TYPE.TITLE, "body": PH_TYPE.BODY, "image1": "Picture Placeholder 3"}
    >>> placeholders = resolve_placeholders(slide, specs)
"""

from __future__ import annotations

import logging
from typing import TYPE_CHECKING, Union

from pptx.enum.shapes import PP_PLACEHOLDER_TYPE as PH_TYPE

if TYPE_CHECKING:
    from pptx.shapes.base import BaseShape
    from pptx.slide import Slide

logger = logging.getLogger(__name__)

# Type alias for placeholder type specification (enum or string)
PlaceholderTypeSpec = Union[PH_TYPE, str]

# Common placeholder type name mappings (string -> enum)
# Allows users to specify types as strings for convenience
_PH_TYPE_ALIASES: dict[str, PH_TYPE] = {
    # Title types
    "title": PH_TYPE.TITLE,
    "TITLE": PH_TYPE.TITLE,
    "center_title": PH_TYPE.CENTER_TITLE,
    "CENTER_TITLE": PH_TYPE.CENTER_TITLE,
    # Subtitle
    "subtitle": PH_TYPE.SUBTITLE,
    "SUBTITLE": PH_TYPE.SUBTITLE,
    # Body/content types
    "body": PH_TYPE.BODY,
    "BODY": PH_TYPE.BODY,
    "content": PH_TYPE.BODY,
    "CONTENT": PH_TYPE.BODY,
    # Picture/image types
    "picture": PH_TYPE.PICTURE,
    "PICTURE": PH_TYPE.PICTURE,
    "image": PH_TYPE.PICTURE,
    "IMAGE": PH_TYPE.PICTURE,
    # Object type
    "object": PH_TYPE.OBJECT,
    "OBJECT": PH_TYPE.OBJECT,
    # Chart
    "chart": PH_TYPE.CHART,
    "CHART": PH_TYPE.CHART,
    # Table
    "table": PH_TYPE.TABLE,
    "TABLE": PH_TYPE.TABLE,
    # Footer elements
    "footer": PH_TYPE.FOOTER,
    "FOOTER": PH_TYPE.FOOTER,
    "date": PH_TYPE.DATE,
    "DATE": PH_TYPE.DATE,
    "slide_number": PH_TYPE.SLIDE_NUMBER,
    "SLIDE_NUMBER": PH_TYPE.SLIDE_NUMBER,
}


def _resolve_ph_type(ph_type: PlaceholderTypeSpec) -> PH_TYPE:
    """Resolve a placeholder type specification to an enum value.
    
    Args:
        ph_type: Either a PH_TYPE enum value or a string alias.
        
    Returns:
        The resolved PH_TYPE enum value.
        
    Raises:
        ValueError: If string alias is not recognized.
    """
    if isinstance(ph_type, PH_TYPE):
        return ph_type
    
    if isinstance(ph_type, str):
        if ph_type in _PH_TYPE_ALIASES:
            return _PH_TYPE_ALIASES[ph_type]
        # Try to match enum by name directly
        try:
            return PH_TYPE[ph_type.upper()]
        except KeyError:
            pass
        raise ValueError(
            f"Unknown placeholder type alias: '{ph_type}'. "
            f"Valid aliases: {sorted(set(_PH_TYPE_ALIASES.keys()))}"
        )
    
    raise TypeError(
        f"ph_type must be PH_TYPE enum or string, got {type(ph_type).__name__}"
    )


def _get_placeholder_info(shape: "BaseShape") -> dict:
    """Extract placeholder information from a shape for diagnostics.
    
    Args:
        shape: A shape from a slide.
        
    Returns:
        Dictionary with placeholder information.
    """
    try:
        ph_format = shape.placeholder_format
        return {
            "name": shape.name,
            "idx": ph_format.idx if ph_format else None,
            "type": ph_format.type if ph_format else None,
            "type_name": ph_format.type.name if ph_format and ph_format.type else None,
        }
    except AttributeError:
        return {"name": shape.name, "idx": None, "type": None, "type_name": None}


def _is_placeholder(shape: "BaseShape") -> bool:
    """Check if a shape is a placeholder.
    
    Args:
        shape: A shape from a slide.
        
    Returns:
        True if the shape is a placeholder.
    """
    try:
        return shape.is_placeholder
    except AttributeError:
        return False


def get_placeholders(
    slide: "Slide",
    name: str | None = None,
    ph_type: PlaceholderTypeSpec | None = None,
) -> list["BaseShape"]:
    """Get all placeholders matching the given criteria.
    
    Finds placeholders by exact name match, placeholder type, or both.
    If both name and ph_type are provided, returns placeholders matching
    either criterion (OR logic).
    
    Args:
        slide: The slide to search for placeholders.
        name: Optional exact name to match against shape.name.
        ph_type: Optional placeholder type (PH_TYPE enum or string alias).
            Common aliases: "title", "body"/"content", "picture"/"image",
            "subtitle", "object", "chart", "table".
            
    Returns:
        List of matching placeholder shapes (may be empty).
        
    Raises:
        ValueError: If ph_type string is not a recognized alias.
        
    Example:
        >>> # Get all title placeholders
        >>> titles = get_placeholders(slide, ph_type="title")
        >>> # Get all placeholders named "Content Placeholder 2"
        >>> content = get_placeholders(slide, name="Content Placeholder 2")
    """
    if name is None and ph_type is None:
        # Return all placeholders if no filter specified
        return [shape for shape in slide.shapes if _is_placeholder(shape)]
    
    resolved_type: PH_TYPE | None = None
    if ph_type is not None:
        resolved_type = _resolve_ph_type(ph_type)
    
    matches = []
    for shape in slide.shapes:
        if not _is_placeholder(shape):
            continue
        
        matched = False
        
        # Check name match
        if name is not None and shape.name == name:
            matched = True
        
        # Check type match
        if resolved_type is not None:
            try:
                ph_format = shape.placeholder_format
                if ph_format and ph_format.type == resolved_type:
                    matched = True
            except AttributeError:
                pass
        
        if matched:
            matches.append(shape)
    
    return matches


def get_placeholder(
    slide: "Slide",
    name: str | None = None,
    ph_type: PlaceholderTypeSpec | None = None,
    raise_on_missing: bool = True,
) -> "BaseShape | None":
    """Get the first placeholder matching the given criteria.
    
    Finds a placeholder by exact name match, placeholder type, or both.
    If multiple placeholders match, logs a warning and returns the first.
    
    Args:
        slide: The slide to search for placeholders.
        name: Optional exact name to match against shape.name.
        ph_type: Optional placeholder type (PH_TYPE enum or string alias).
            Common aliases: "title", "body"/"content", "picture"/"image",
            "subtitle", "object", "chart", "table".
        raise_on_missing: If True (default), raises PlaceholderNotFoundError
            when no matching placeholder is found. If False, returns None.
            
    Returns:
        The first matching placeholder shape, or None if not found and
        raise_on_missing is False.
        
    Raises:
        PlaceholderNotFoundError: If no matching placeholder found and
            raise_on_missing is True.
        ValueError: If both name and ph_type are None.
        ValueError: If ph_type string is not a recognized alias.
        
    Example:
        >>> from pptx.enum.shapes import PP_PLACEHOLDER_TYPE as PH_TYPE
        >>> 
        >>> # Get title placeholder by type
        >>> title = get_placeholder(slide, ph_type=PH_TYPE.TITLE)
        >>> title.text_frame.text = "My Title"
        >>> 
        >>> # Get by name
        >>> content = get_placeholder(slide, name="Content Placeholder 2")
        >>> 
        >>> # Return None instead of raising
        >>> maybe_subtitle = get_placeholder(slide, ph_type="subtitle", raise_on_missing=False)
        >>> if maybe_subtitle:
        ...     maybe_subtitle.text_frame.text = "Optional subtitle"
    """
    if name is None and ph_type is None:
        raise ValueError("At least one of 'name' or 'ph_type' must be provided")
    
    matches = get_placeholders(slide, name=name, ph_type=ph_type)
    
    if not matches:
        if raise_on_missing:
            # Build diagnostic message
            all_placeholders = get_placeholders(slide)
            available_info = [_get_placeholder_info(p) for p in all_placeholders]
            
            search_criteria = []
            if name:
                search_criteria.append(f"name='{name}'")
            if ph_type:
                if isinstance(ph_type, PH_TYPE):
                    search_criteria.append(f"type={ph_type.name}")
                else:
                    search_criteria.append(f"type='{ph_type}'")
            
            available_str = "\n".join(
                f"  - name='{p['name']}', type={p['type_name']}, idx={p['idx']}"
                for p in available_info
            )
            
            raise PlaceholderNotFoundError(
                f"No placeholder found matching {', '.join(search_criteria)}.\n"
                f"Available placeholders on this slide:\n{available_str}"
            )
        return None
    
    if len(matches) > 1:
        search_criteria = []
        if name:
            search_criteria.append(f"name='{name}'")
        if ph_type:
            if isinstance(ph_type, PH_TYPE):
                search_criteria.append(f"type={ph_type.name}")
            else:
                search_criteria.append(f"type='{ph_type}'")
        
        match_names = [m.name for m in matches]
        logger.warning(
            f"Multiple placeholders match {', '.join(search_criteria)}: "
            f"{match_names}. Using first: '{matches[0].name}'"
        )
    
    return matches[0]


def resolve_placeholders(
    slide: "Slide",
    specs: dict[str, PlaceholderTypeSpec | str],
) -> dict[str, "BaseShape"]:
    """Resolve multiple placeholder specifications at once.
    
    Takes a dictionary mapping field names to placeholder specifications
    (either placeholder type enums/strings or exact placeholder names)
    and returns a dictionary mapping those field names to the resolved
    placeholder shapes.
    
    The specification values are interpreted as follows:
    - PH_TYPE enum: Match by placeholder type
    - String in _PH_TYPE_ALIASES (e.g., "title", "body"): Match by type 
    - Other string: Match by exact placeholder name
    
    Args:
        slide: The slide to search for placeholders.
        specs: Dictionary mapping field names to specifications.
            Keys are arbitrary field names (e.g., "title", "content", "image1").
            Values are either:
            - PH_TYPE enum values (e.g., PH_TYPE.TITLE)
            - Type alias strings (e.g., "title", "body", "picture")
            - Exact placeholder names (e.g., "Picture Placeholder 3")
            
    Returns:
        Dictionary mapping field names to resolved placeholder shapes.
        Only includes entries where placeholders were found.
        
    Raises:
        PlaceholderNotFoundError: If any specification cannot be resolved.
        
    Example:
        >>> from pptx.enum.shapes import PP_PLACEHOLDER_TYPE as PH_TYPE
        >>> 
        >>> specs = {
        ...     "title": PH_TYPE.TITLE,
        ...     "content": "body",  # alias for PH_TYPE.BODY
        ...     "image1": "Picture Placeholder 3",  # exact name
        ... }
        >>> resolved = resolve_placeholders(slide, specs)
        >>> resolved["title"].text_frame.text = "My Title"
        >>> resolved["content"].text_frame.text = "Body text here"
    """
    resolved: dict[str, "BaseShape"] = {}
    
    for field_name, spec in specs.items():
        # Determine if spec is a type or a name
        is_type_spec = False
        
        if isinstance(spec, PH_TYPE):
            is_type_spec = True
        elif isinstance(spec, str):
            # Check if it's a known type alias
            if spec in _PH_TYPE_ALIASES:
                is_type_spec = True
            else:
                # Try enum name lookup
                try:
                    PH_TYPE[spec.upper()]
                    is_type_spec = True
                except KeyError:
                    pass
        
        try:
            if is_type_spec:
                placeholder = get_placeholder(slide, ph_type=spec, raise_on_missing=True)
            else:
                # Treat as exact name match
                placeholder = get_placeholder(slide, name=str(spec), raise_on_missing=True)
            
            if placeholder is not None:
                resolved[field_name] = placeholder
                
        except PlaceholderNotFoundError as e:
            # Re-raise with field context
            raise PlaceholderNotFoundError(
                f"Cannot resolve field '{field_name}': {e}"
            ) from e
    
    return resolved


class PlaceholderNotFoundError(Exception):
    """Raised when a placeholder cannot be found on a slide."""
    pass
