"""Unified slide construction functionality.

This module provides generic functions for building and populating slides
using the layout registry and placeholder resolver.
"""

from __future__ import annotations

import re
import logging
from typing import TYPE_CHECKING

from pptx.util import Pt
from pptx.dml.color import RGBColor

from .config import Config
from pathlib import Path

from .images import add_images_for_layout, add_background_image
from .layout_discovery import LayoutRegistry, validate_layout_name
from .markdown_parser import SlideData, SPACER_MARKER
from .placeholder_resolver import resolve_placeholders, get_placeholders, PlaceholderNotFoundError
from .rich_text import add_formatted_text, add_bullet, remove_bullet, add_numbering

if TYPE_CHECKING:
    from pptx.presentation import Presentation
    from pptx.slide import Slide
    from pptx.text.text import TextFrame

logger = logging.getLogger(__name__)

# Tag used to mark which layout TextBox was used for a given semantic prompt
# (title/subtitle/body/etc.). We use this because we may clear the layout
# "Click to add …" text to prevent it from rendering behind a cloned
# slide-level textbox. After clearing, text-based matching no longer works.
_LAYOUT_PROMPT_TAG_PREFIX = "iltci_prompt:"


def _get_layout_shape_descr(shape) -> str:
    """Best-effort read of a layout shape's non-visual description field."""
    try:
        cNvPr = shape._element.nvSpPr.cNvPr
        return cNvPr.get("descr") or ""
    except Exception:
        return ""


def _append_layout_shape_descr_tag(shape, tag: str) -> None:
    """Append a tag to the layout shape's non-visual description field."""
    try:
        cNvPr = shape._element.nvSpPr.cNvPr
        existing = (cNvPr.get("descr") or "").strip()
        if tag in existing:
            return
        new_descr = f"{existing} {tag}".strip() if existing else tag
        cNvPr.set("descr", new_descr)
    except Exception:
        # Non-fatal; tagging is only an optimization for later lookup.
        return


def build_slide(layout_name: str, prs: "Presentation", registry: LayoutRegistry) -> "Slide":
    """Build a slide using the specified layout from the registry.
    
    This is the primary function for creating new slides. It validates the
    layout name against the registry, retrieves the layout index, and creates
    a new slide.
    
    Args:
        layout_name: Exact layout name from the template (must exist in registry).
        prs: PowerPoint presentation object.
        registry: LayoutRegistry mapping layout names to indices.
        
    Returns:
        Created slide object.
        
    Raises:
        ValueError: If layout_name is not found in the registry.
        
    Example:
        >>> registry = load_layout_registry(prs)
        >>> slide = build_slide("Title and Content", prs, registry)
    """
    # Validate layout name against registry
    validate_layout_name(layout_name, registry, raise_on_missing=True)
    
    # Get layout index and create slide
    layout_idx = registry[layout_name]
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)
    
    logger.info(f"Built slide using layout '{layout_name}' (index {layout_idx})")
    logger.debug(f"  Slide has {len(slide.shapes)} shapes")
    
    return slide


def _find_layout_shape_by_prompt(slide: "Slide", prompt_keyword: str):
    """Find a TextBox shape on the slide's layout whose default text contains *prompt_keyword*.

    Many modern templates use free-form TextBox shapes (not typed
    placeholders) on the layout.  These shapes show "Click to add …"
    prompt text that is visible in design mode.  When a slide is created
    from such a layout the shapes are *not* inherited as editable
    slide-level shapes; they render from the layout behind the slide
    content.

    To "fill" such a shape we need to:
    1. Locate it on the layout by matching its prompt text.
    2. Clone the shape's XML to the slide (preserving position, size,
       and formatting).
    3. Set the desired text on the clone.

    This helper performs step 1.

    Args:
        slide: The slide whose *layout* will be searched.
        prompt_keyword: Case-insensitive substring to match (e.g.
            ``"title"``, ``"subtitle"``, ``"body"``).

    Returns:
        The first matching layout shape, or ``None``.
    """
    keyword_lower = prompt_keyword.lower()

    # 1) Prefer a previously-tagged shape. This allows us to keep finding the
    # correct layout shape even after we've cleared its prompt text.
    tag = f"{_LAYOUT_PROMPT_TAG_PREFIX}{keyword_lower}"
    for shape in slide.slide_layout.shapes:
        descr = _get_layout_shape_descr(shape).lower()
        if tag in descr:
            return shape

    # 2) Fallback: match by prompt text content (e.g., "Click to add title")
    #    OR by shape name (e.g., "ph_feature_desc_1").
    for shape in slide.slide_layout.shapes:
        if shape.has_text_frame:
            shape_text = shape.text_frame.text.lower()
            shape_name = getattr(shape, "name", "").lower()
            if keyword_lower in shape_text or keyword_lower in shape_name:
                # Tag for future lookups (important if we clear the prompt text)
                _append_layout_shape_descr_tag(shape, tag)
                return shape
    return None


# ---------------------------------------------------------------------------
# Keywords used by the fallback content-shape search in populate_slide.
# ---------------------------------------------------------------------------
_CONTENT_KEYWORDS: list[str] = ["content", "body", "desc", "feature"]
_TITLE_KEYWORDS: list[str] = ["title"]


def _find_all_layout_shapes_by_keywords(
    slide: "Slide",
    keywords: list[str],
    *,
    already_used: set | None = None,
) -> list:
    """Find *all* layout TextBox shapes matching any of *keywords*.

    Searches both ``shape.name`` and ``shape.text_frame.text`` (case-
    insensitive).  Previously-tagged shapes (via the ``descr`` tag) are
    also considered.  Shapes already consumed by an earlier call can be
    excluded via *already_used* (a set of ``shape.shape_id``).

    Returns:
        A list of matching layout shapes in document order (preserves the
        order they appear in the layout's shape tree).
    """
    used = already_used or set()
    matches: list = []
    kw_lowers = [kw.lower() for kw in keywords]

    for shape in slide.slide_layout.shapes:
        if shape.shape_id in used:
            continue

        # Check tag first (previously matched & tagged shapes).
        descr = _get_layout_shape_descr(shape).lower()
        for kw in kw_lowers:
            tag = f"{_LAYOUT_PROMPT_TAG_PREFIX}{kw}"
            if tag in descr:
                matches.append(shape)
                used.add(shape.shape_id)
                break
        else:
            # No tag matched – try name and text.
            if not getattr(shape, "has_text_frame", False):
                continue
            shape_text = shape.text_frame.text.lower()
            shape_name = getattr(shape, "name", "").lower()
            for kw in kw_lowers:
                if kw in shape_text or kw in shape_name:
                    # Tag it for future lookups.
                    tag = f"{_LAYOUT_PROMPT_TAG_PREFIX}{kw}"
                    _append_layout_shape_descr_tag(shape, tag)
                    matches.append(shape)
                    used.add(shape.shape_id)
                    break

    return matches


# Cache of original layout shape XML elements.
# When a layout TextBox is cloned to a slide, the layout copy is cleared to
# prevent visual overlap.  Subsequent slides sharing the same layout would
# clone from the *cleared* element and lose all formatting.  This cache
# stores a pristine deep-copy of the element **before** the first clear so
# that every clone preserves the original font, size, color, etc.
# Key: ``(id(slide_layout_object), shape_id)`` — both are stable across
# python-pptx wrapper re-creations within a single generator run.
_layout_shape_originals: dict[tuple, object] = {}


def _shape_cache_key(layout_shape) -> tuple:
    """Build a stable cache key for a layout shape.

    ``id(element)`` is NOT reliable because lxml may return different
    proxy objects for the same underlying node.  Instead we combine the
    Python ``id`` of the *SlideLayout* (which **is** a stable singleton
    per layout within a Presentation) with the shape's numeric ``id``
    attribute from the XML (``<p:cNvPr id="…">``).
    """
    try:
        # layout_shape lives inside a slide layout's spTree; walk up to
        # obtain the layout wrapper.  python-pptx shapes don't expose this
        # directly, so we rely on the part → package → Presentation chain.
        # Easier: get shape_id (== cNvPr @id) which is unique per layout.
        return (id(layout_shape.part), layout_shape.shape_id)
    except Exception:
        # Fallback — should not happen in normal operation.
        return (id(layout_shape), getattr(layout_shape, "shape_id", 0))


def _clone_layout_shape_with_text(slide: "Slide", layout_shape, text: str) -> bool:
    """Clone a layout shape to the slide and set its text.

    Creates a deep copy of *layout_shape*'s XML element and appends it
    to the slide's shape tree.  The first paragraph of the cloned shape
    is then set to *text*, preserving the original run-level formatting
    (font family, size, color, etc.) from the layout.

    To ensure formatting survives across multiple slides that share the
    same layout, the original element XML is cached before the layout
    shape is cleared.  Subsequent clones are made from this cache.

    Args:
        slide: Target slide.
        layout_shape: Shape from the slide layout to clone.
        text: Text to set on the cloned shape.

    Returns:
        ``True`` if the clone was created and text set, ``False`` on error.
    """
    from copy import deepcopy

    try:
        # NOTE: Many templates in this project use free-form TextBox shapes on the
        # *layout* with literal prompt text like "Click to add title".
        # When we clone a layout TextBox onto the slide to populate it, the
        # original layout TextBox still renders behind the slide content unless
        # we also clear its prompt text.
        layout_prompt_text = ""
        if getattr(layout_shape, "has_text_frame", False):
            try:
                layout_prompt_text = layout_shape.text_frame.text or ""
            except Exception:
                layout_prompt_text = ""

        # Cache the pristine element before any modification so that
        # subsequent slides cloning the same layout shape get full formatting.
        shape_key = _shape_cache_key(layout_shape)
        if shape_key not in _layout_shape_originals:
            _layout_shape_originals[shape_key] = deepcopy(layout_shape._element)

        cloned_sp = deepcopy(_layout_shape_originals[shape_key])
        slide.shapes._spTree.append(cloned_sp)

        # Walk the cloned shape to set text while preserving formatting.
        # We locate the cloned shape by finding the element we just appended.
        cloned_shape = None
        for s in slide.shapes:
            if s._element is cloned_sp:
                cloned_shape = s
                break

        if cloned_shape is not None and cloned_shape.has_text_frame:
            tf = cloned_shape.text_frame
            # Preserve formatting of first paragraph / first run
            if tf.paragraphs:
                para = tf.paragraphs[0]
                if para.runs:
                    para.runs[0].text = text
                    # Remove extra runs (prompt text spans)
                    for run in para.runs[1:]:
                        run.text = ""
                else:
                    para.text = text
                # Remove subsequent paragraphs (some prompts span multiple)
                for extra_para in tf.paragraphs[1:]:
                    extra_para.text = ""
            else:
                tf.text = text

            # Clear the underlying layout prompt so it doesn't visually overlap
            # with the cloned slide-level textbox.  Any non-empty layout text is
            # treated as a prompt—not just "Click to add …" strings.
            if (
                getattr(layout_shape, "has_text_frame", False)
                and layout_prompt_text
            ):
                try:
                    layout_shape.text_frame.clear()
                    # Ensure no stray text remains in the default paragraph
                    layout_shape.text_frame.text = ""
                    logger.debug(
                        "Cleared layout prompt text on shape '%s' (was %r)",
                        getattr(layout_shape, "name", "<unnamed>"),
                        layout_prompt_text,
                    )
                except Exception as exc:
                    logger.warning(
                        "Failed to clear layout prompt text on shape '%s': %s",
                        getattr(layout_shape, "name", "<unnamed>"),
                        exc,
                    )
            return True
    except Exception as exc:
        logger.warning(f"Failed to clone layout shape: {exc}")
    return False


# Frontmatter keys that are handled explicitly or are metadata—not layout text fields.
_RESERVED_FM_KEYS: frozenset[str] = frozenset({
    "layout", "id", "title", "subtitle", "body", "background",
    "images", "section_name", "_normalized_images", "marp", "theme",
    "author", "date",
})


def apply_frontmatter_to_slide(
    slide: "Slide",
    data: SlideData,
    config: Config,
) -> None:
    """Apply frontmatter-driven actions to a slide.

    Handles semantic keys from the slide's ``frontmatter`` dict:

    - **title** / **subtitle** / **body**: populated via typed
      placeholders first; falls back to cloning matching TextBox shapes
      from the layout (for templates that use free-form shapes).
    - **background.image_path**: resolves the path relative to
      :pyattr:`Config.assets_dir` and calls :func:`add_background_image`.
    - **Other keys**: any remaining string-valued frontmatter key is
      treated as a layout TextBox search term (underscores → spaces).
      For example ``section_label: "Live Demo"`` will find a layout
      shape whose text contains "section label" and clone it with the
      provided value.

    Args:
        slide: PowerPoint slide object.
        data: SlideData with populated ``frontmatter`` dict.
        config: Configuration object (used for ``assets_dir``).
    """
    fm = data.frontmatter
    if not fm:
        return

    slide_label = data.slide_id or data.title or data.layout_name

    # --- background.image_path ---
    bg = fm.get('background')
    if isinstance(bg, dict):
        bg_image_path = bg.get('image_path')
        if bg_image_path:
            resolved = Path(config.assets_dir) / bg_image_path
            logger.info(f"  [{slide_label}] Applying background image: {resolved}")
            add_background_image(slide, resolved)

    # --- title ---
    title_text = fm.get('title')
    if title_text:
        title_set = False
        # Try typed placeholders first
        for title_type in ["title", "center_title"]:
            try:
                result = resolve_placeholders(slide, {"title": title_type})
                ph = result.get("title")
                if ph:
                    ph.text_frame.text = str(title_text)
                    title_set = True
                    logger.info(f"  [{slide_label}] Set title via placeholder ({title_type})")
                    break
            except PlaceholderNotFoundError:
                continue

        # Fallback: find a BODY placeholder whose layout shape name contains
        # "title" (but not "feature_title" etc.) and populate it directly.
        # This handles layouts where the title placeholder has BODY type
        # instead of TITLE type (e.g. Feature Grid).
        if not title_set:
            for lph in slide.slide_layout.placeholders:
                lname = getattr(lph, "name", "").lower()
                if "title" in lname and "feature" not in lname:
                    ph_idx = lph.placeholder_format.idx
                    try:
                        slide_ph = slide.placeholders[ph_idx]
                        slide_ph.text_frame.text = str(title_text)
                        title_set = True
                        logger.info(
                            "  [%s] Set title via BODY placeholder "
                            "idx=%d (layout: '%s')",
                            slide_label, ph_idx, lph.name,
                        )
                        break
                    except (KeyError, AttributeError):
                        pass

        # Final fallback: clone layout TextBox matching "title"
        if not title_set:
            layout_shape = _find_layout_shape_by_prompt(slide, "title")
            if layout_shape:
                if _clone_layout_shape_with_text(slide, layout_shape, str(title_text)):
                    title_set = True
                    logger.info(f"  [{slide_label}] Set title via layout TextBox clone")
            if not title_set:
                logger.warning(f"  [{slide_label}] Could not find title shape")

    # --- subtitle ---
    subtitle_text = fm.get('subtitle')
    if subtitle_text:
        sub_set = False
        # Try typed placeholder first
        try:
            result = resolve_placeholders(slide, {"subtitle": "subtitle"})
            ph = result.get("subtitle")
            if ph:
                ph.text_frame.text = str(subtitle_text)
                sub_set = True
                logger.info(f"  [{slide_label}] Set subtitle via placeholder")
        except PlaceholderNotFoundError:
            pass

        # Fallback: clone layout TextBox matching "subtitle" or "contact"
        if not sub_set:
            for kw in ("subtitle", "contact"):
                layout_shape = _find_layout_shape_by_prompt(slide, kw)
                if layout_shape:
                    if _clone_layout_shape_with_text(slide, layout_shape, str(subtitle_text)):
                        sub_set = True
                        logger.info(f"  [{slide_label}] Set subtitle via layout TextBox clone ('{kw}')")
                        break
            if not sub_set:
                logger.debug(f"  [{slide_label}] No subtitle shape found; skipping")

    # --- body (from frontmatter only, markdown body handled by populate_slide) ---
    body_text = fm.get('body')
    if body_text and not data.content_blocks:
        body_set = False
        for body_type in ["body", "object"]:
            try:
                result = resolve_placeholders(slide, {"body": body_type})
                ph = result.get("body")
                if ph:
                    ph.text_frame.text = str(body_text).strip()
                    body_set = True
                    logger.info(f"  [{slide_label}] Set body via placeholder ({body_type})")
                    break
            except PlaceholderNotFoundError:
                continue

        if not body_set:
            layout_shape = _find_layout_shape_by_prompt(slide, "body")
            if layout_shape:
                if _clone_layout_shape_with_text(slide, layout_shape, str(body_text).strip()):
                    body_set = True
                    logger.info(f"  [{slide_label}] Set body via layout TextBox clone")

    # --- Generic layout TextBox fields from remaining frontmatter keys ---
    # Any string-valued key not already handled above is treated as a layout
    # TextBox search term (underscores converted to spaces).  For example,
    # ``section_label: "Live Demo"`` finds the layout shape whose text
    # contains "section label" and clones it with the provided value.
    for key, value in fm.items():
        if key in _RESERVED_FM_KEYS or not isinstance(value, str):
            continue
        search_term = key.replace("_", " ")
        layout_shape = _find_layout_shape_by_prompt(slide, search_term)
        if layout_shape:
            if _clone_layout_shape_with_text(slide, layout_shape, value):
                logger.info(
                    f"  [{slide_label}] Set '{key}' via layout TextBox clone"
                )
            else:
                logger.warning(
                    f"  [{slide_label}] Found layout shape for '{key}' but clone failed"
                )
        else:
            logger.debug(
                f"  [{slide_label}] No layout shape matching '{search_term}' for key '{key}'"
            )


def _clear_unused_layout_text(slide: "Slide") -> None:
    """Clear text on layout TextBox shapes that were not populated.

    After slide population, any layout TextBox shape whose text has not
    been cloned to the slide will still render behind the slide content.
    This function blanks out any remaining non-empty layout TextBox text
    so that unpopulated prompt text (e.g. "SECTION LABEL", "Feature 1")
    does not leak into the final presentation.

    Shapes that were already handled during population are identified by
    their ``_LAYOUT_PROMPT_TAG_PREFIX`` tag and skipped (they have
    already been cleared by :func:`_clone_layout_shape_with_text`).

    .. note::
       Clearing layout shape text mutates the shared layout object, so
       the change is visible to **all** subsequent slides using the same
       layout.  This is acceptable because the generator either populates
       a shape via frontmatter (in which case the layout text is already
       cleared by the clone step) or the user intentionally left the
       field undefined (in which case blank is the desired output).
    """
    for shape in slide.slide_layout.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        # Skip shapes already handled (tagged during clone operations)
        descr = _get_layout_shape_descr(shape)
        if _LAYOUT_PROMPT_TAG_PREFIX in descr:
            continue
        # Cache original XML before clearing so that future slides sharing
        # this layout can still clone with full formatting preserved.
        from copy import deepcopy                         # noqa: E402 (local)
        shape_key = _shape_cache_key(shape)
        if shape_key not in _layout_shape_originals:
            _layout_shape_originals[shape_key] = deepcopy(shape._element)
        # Clear the remaining prompt text
        try:
            original_text = shape.text_frame.text
            shape.text_frame.clear()
            shape.text_frame.text = ""
            logger.debug(
                "Cleared unused layout text on shape '%s' (was %r)",
                getattr(shape, "name", "<unnamed>"),
                original_text,
            )
        except Exception as exc:
            logger.warning(
                "Failed to clear unused layout text on shape '%s': %s",
                getattr(shape, "name", "<unnamed>"),
                exc,
            )


def populate_slide(
    slide: "Slide",
    data: SlideData,
    config: Config,
    registry: LayoutRegistry | None = None,
) -> None:
    """Populate a slide with content from SlideData.
    
    Uses the placeholder resolver to find title and content placeholders,
    then populates them with the provided data. If images are present and
    a registry is provided, uses add_images_for_layout for config-driven
    image placement.
    
    Args:
        slide: Slide object to populate.
        data: SlideData containing title, content_blocks, images, etc.
        config: Configuration object for fonts and formatting.
        registry: Optional LayoutRegistry for image placement validation.
        
    Raises:
        PlaceholderNotFoundError: If required placeholders cannot be found.
        
    Example:
        >>> slide = build_slide("Title and Content", prs, registry)
        >>> populate_slide(slide, slide_data, config, registry)
    """
    # Apply frontmatter-driven actions (background, subtitle, etc.)
    apply_frontmatter_to_slide(slide, data, config)
    
    # Handle images using PICTURE placeholders from the template
    if data.images:
        if registry is None:
            logger.warning(
                f"Slide '{data.title or data.layout_name}' has {len(data.images)} images "
                "but no registry provided. Images will not be placed."
            )
        else:
            add_images_for_layout(data, slide, config, registry)
    
    # --- Title (only if NOT already set by apply_frontmatter_to_slide) ---
    # Frontmatter-based title is handled in apply_frontmatter_to_slide.
    # Here we handle title from content H1 (when no frontmatter title exists).
    fm_has_title = data.frontmatter and data.frontmatter.get('title')
    title_ph = None
    content_ph = None
    
    if data.title and not fm_has_title:
        for title_type in ["title", "center_title"]:
            try:
                result = resolve_placeholders(slide, {"title": title_type})
                title_ph = result.get("title")
                if title_ph:
                    logger.debug(f"  Found title placeholder with type '{title_type}'")
                    break
            except PlaceholderNotFoundError:
                continue
        
        # Fallback: clone layout TextBox for title
        if title_ph is None:
            layout_shape = _find_layout_shape_by_prompt(slide, "title")
            if layout_shape:
                _clone_layout_shape_with_text(slide, layout_shape, data.title)
                logger.debug(f"  Set title via layout TextBox clone: '{data.title}'")
            else:
                logger.warning(f"  No title placeholder found for slide '{data.layout_name}'")
        else:
            title_ph.text_frame.text = data.title
            logger.debug(f"  Set title: '{data.title}'")
    
    # --- Content blocks ---
    if data.content_blocks:
        _multi_body_done = False
        for content_type in ["body", "object", "subtitle"]:
            try:
                if content_type == "body":
                    body_matches = get_placeholders(slide, ph_type="body")
                    if len(body_matches) > 2:
                        # Multi-BODY layout (e.g. Feature Grid with
                        # ph_feature_title_1…6, ph_feature_desc_1…6):
                        # populate slide placeholders directly by idx order
                        # instead of cloning layout shapes.
                        sorted_phs = sorted(
                            body_matches,
                            key=lambda ph: ph.placeholder_format.idx,
                        )
                        # Identify the title placeholder idx to exclude
                        # from content distribution.
                        title_idx = None
                        if data.title:
                            for lph in slide.slide_layout.placeholders:
                                lname = getattr(lph, "name", "").lower()
                                if "title" in lname and "feature" not in lname:
                                    title_idx = lph.placeholder_format.idx
                                    break
                        content_phs = [
                            ph for ph in sorted_phs
                            if ph.placeholder_format.idx != title_idx
                        ]
                        for i, ph in enumerate(content_phs):
                            if i < len(data.content_blocks):
                                build_rich_content(
                                    ph.text_frame,
                                    [data.content_blocks[i]],
                                    config,
                                    data.layout_name,
                                )
                                logger.debug(
                                    "  Multi-BODY: content_block[%d] → "
                                    "placeholder idx=%d (%s)",
                                    i, ph.placeholder_format.idx, ph.name,
                                )
                            else:
                                ph.text_frame.clear()
                                ph.text_frame.text = ""
                        logger.debug(
                            "  Multi-BODY: populated %d/%d placeholders "
                            "with %d content blocks",
                            min(len(content_phs), len(data.content_blocks)),
                            len(content_phs),
                            len(data.content_blocks),
                        )
                        _multi_body_done = True
                        break
                result = resolve_placeholders(slide, {"content": content_type})
                content_ph = result.get("content")
                if content_ph:
                    logger.debug(f"  Found content placeholder with type '{content_type}'")
                    break
            except PlaceholderNotFoundError:
                continue
        
        if _multi_body_done:
            pass  # Content already distributed to individual placeholders
        elif content_ph is None:
            # Fallback: find layout TextBox shapes matching content keywords.
            # Collect *all* matching shapes so that multi-shape layouts (e.g.
            # 'Feature Grid' with ph_feature_desc_1 … ph_feature_desc_N) can
            # each receive a portion of the content blocks.
            content_shapes = _find_all_layout_shapes_by_keywords(
                slide, _CONTENT_KEYWORDS,
            )
            if content_shapes:
                n_shapes = len(content_shapes)
                n_blocks = len(data.content_blocks)
                # Distribute content blocks as evenly as possible across shapes.
                base, extra = divmod(n_blocks, n_shapes)
                idx = 0
                for i, layout_shape in enumerate(content_shapes):
                    count = base + (1 if i < extra else 0)
                    block_slice = data.content_blocks[idx : idx + count]
                    idx += count
                    if not block_slice:
                        continue
                    if _clone_layout_shape_with_text(slide, layout_shape, ""):
                        cloned = slide.shapes[-1] if slide.shapes else None
                        if cloned and cloned.has_text_frame:
                            build_rich_content(
                                cloned.text_frame,
                                block_slice,
                                config,
                                data.layout_name,
                            )
                            logger.debug(
                                "  Set content via layout TextBox clone "
                                "(shape %d/%d, %d blocks)",
                                i + 1,
                                n_shapes,
                                len(block_slice),
                            )
            else:
                logger.warning(f"  No content placeholder found for slide '{data.layout_name}'")
        else:
            text_frame = content_ph.text_frame
            build_rich_content(text_frame, data.content_blocks, config, data.layout_name)
            logger.debug(f"  Populated {len(data.content_blocks)} content blocks")

    # --- Cleanup: blank out any remaining layout text that wasn't populated ---
    _clear_unused_layout_text(slide)


def _normalize_layout_key(layout_name: str) -> str:
    """Normalize layout name to a config-friendly key.
    
    Converts layout names like "Image Right" to "image_right" for config lookup.
    
    Args:
        layout_name: Layout name from template (e.g., "Image Right", "Dual Image")
        
    Returns:
        Normalized key (e.g., "image_right", "dual_image")
    """
    return layout_name.lower().replace(" ", "_")


def _extract_template_props(text_frame: "TextFrame") -> dict:
    """Extract formatting properties from a text frame before clearing.

    Captures font name, color, size, and line spacing from the first
    paragraph/run so they can be re-applied as fallback defaults after
    ``text_frame.clear()`` destroys the template formatting.

    Args:
        text_frame: PowerPoint text frame to extract properties from.

    Returns:
        Dict with optional keys: ``font_name`` (str), ``font_color``
        (:class:`~pptx.util.RGBColor`), ``font_size`` (EMU value),
        ``line_spacing`` (float or :class:`~pptx.util.Pt`).
    """
    props: dict = {}
    try:
        if text_frame.paragraphs:
            para = text_frame.paragraphs[0]
            if para.line_spacing is not None:
                props['line_spacing'] = para.line_spacing
            if para.runs:
                run = para.runs[0]
                if run.font.name:
                    props['font_name'] = run.font.name
                if run.font.size is not None:
                    props['font_size'] = run.font.size
                try:
                    rgb = run.font.color.rgb
                    if rgb is not None:
                        props['font_color'] = rgb
                except (AttributeError, TypeError):
                    pass
    except Exception:
        logger.debug("Could not extract template properties from text frame")
    return props


def build_rich_content(
    text_frame: "TextFrame",
    content_blocks: list[str],
    config: Config,
    layout_name: str = "",
) -> None:
    """Build rich content from content blocks into a text frame.
    
    Processes content blocks (paragraphs, bullets, headers, spacers) and
    applies appropriate formatting based on configuration.
    
    Font sizes are looked up in order:
    1. fonts.<layout_key>.<property>_pt (e.g., fonts.text.h2_header_pt)
    2. fonts.<layout_key>.<property> (e.g., fonts.text.h2_header)
    3. Default values
    
    Template formatting (font name, color, line spacing) is extracted from
    the text frame before clearing and re-applied as fallbacks when no
    config override is present.
    
    Supports:
    - Headers: ## (H2), ### (H3), #### (H4), ##### (H5)
    - Bullets: - (level 0), "  - " (level 1)
    - Numbered lists: 1. 2. 3. etc.
    - Spacers: SPACER_MARKER for vertical spacing
    - Plain text with markdown formatting (**bold**, *italic*, etc.)
    
    Args:
        text_frame: PowerPoint text frame to populate.
        content_blocks: List of content strings to process.
        config: Configuration object for fonts and formatting.
        layout_name: Layout name for config lookup (e.g., "Text", "Image Right")
        
    Example:
        >>> build_rich_content(text_frame, ["## Introduction", "- Point 1"], config, "Text")
    """
    # Extract template formatting before clearing (clear destroys it)
    template_props = _extract_template_props(text_frame)
    
    # Clear existing content
    text_frame.clear()
    
    # Normalize layout name for config lookup
    layout_key = _normalize_layout_key(layout_name) if layout_name else ""
    
    def get_font_size(prop: str, default: int) -> int:
        """Get font size from config, trying layout-specific then fallback paths."""
        if layout_key:
            # Try layout-specific with _pt suffix first
            val = config.get(f"fonts.{layout_key}.{prop}_pt", None)
            if val is not None:
                return val
            # Try layout-specific without suffix
            val = config.get(f"fonts.{layout_key}.{prop}", None)
            if val is not None:
                return val
        return default
    
    def get_font_prop(prop: str, default=None):
        """Get a font property from config, trying layout-specific then global.
        
        Used for non-size font properties (font_name, font_color) where the
        template value serves as the ultimate fallback.
        """
        if layout_key:
            val = config.get(f"fonts.{layout_key}.{prop}", None)
            if val is not None:
                return val
        val = config.get(f"fonts.{prop}", None)
        if val is not None:
            return val
        return default
    
    def get_formatting(prop: str, default: bool) -> bool:
        """Get formatting setting from config, trying layout-specific then fallback."""
        if layout_key:
            # Try layout-specific formatting
            val = config.get(f"formatting.{layout_key}.{prop}", None)
            if val is not None:
                return val
        # Try global formatting
        val = config.get(f"formatting.{prop}", None)
        if val is not None:
            return val
        return default
    
    def get_spacing(prop: str, default):
        """Get spacing setting from config, trying layout-specific then fallback."""
        if layout_key:
            # Try layout-specific spacing
            val = config.get(f"spacing.{layout_key}.{prop}", None)
            if val is not None:
                return val
        # Try global spacing
        val = config.get(f"spacing.{prop}", None)
        if val is not None:
            return val
        return default
    
    # Get font sizes using layout-aware lookup
    # Use template font_size (EMU) converted to pt as fallback for body-level defaults
    _tpl_font_size_emu = template_props.get('font_size')
    _tpl_fs = round(_tpl_font_size_emu / 12700) if _tpl_font_size_emu else 24
    h2_size = get_font_size("h2_header", _tpl_fs + 4)
    h3_size = get_font_size("h3_header", _tpl_fs)
    h4_size = get_font_size("h4_header", max(_tpl_fs - 4, 10))
    h5_size = get_font_size("h5_header", max(_tpl_fs - 6, 8))
    body_size = get_font_size("body_text", _tpl_fs)
    bullet_size = get_font_size("bullet", _tpl_fs)
    numbered_size = get_font_size("numbered", _tpl_fs)
    spacer_size = get_font_size("spacer", max(_tpl_fs // 2, 6))
    numbering_type = config.get("bullets.numbering_type", "arabicPeriod")
    
    # Get bold settings using layout-aware lookup
    h2_bold = get_formatting("h2_bold", True)
    h3_bold = get_formatting("h3_bold", False)
    h4_bold = get_formatting("h4_bold", False)
    h5_bold = get_formatting("h5_bold", False)
    
    # Get spacing settings using layout-aware lookup
    # Use template line_spacing as fallback when config doesn't specify
    line_spacing = get_spacing("line_spacing", template_props.get('line_spacing'))
    space_after_pt = get_spacing("space_after_pt", None)  # None = use template default
    
    # Resolve font name: config overrides > template default
    font_name = get_font_prop("font_name", template_props.get('font_name'))
    
    # Resolve font color: config (hex string) > template (RGBColor)
    font_color_cfg = get_font_prop("font_color", None)
    if font_color_cfg is not None:
        hex_str = str(font_color_cfg).lstrip('#')
        font_color: RGBColor | None = RGBColor(
            int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16)
        )
    else:
        font_color = template_props.get('font_color')
    
    # Process each content block
    for block in content_blocks:
        # Each block may contain multiple lines
        for line in block.split("\n"):
            line_stripped = line.strip()
            if not line_stripped:
                continue
            
            _add_content_line(
                text_frame,
                line_stripped,
                h2_size=h2_size,
                h3_size=h3_size,
                h4_size=h4_size,
                h5_size=h5_size,
                body_size=body_size,
                bullet_size=bullet_size,
                numbered_size=numbered_size,
                spacer_size=spacer_size,
                numbering_type=numbering_type,
                h2_bold=h2_bold,
                h3_bold=h3_bold,
                h4_bold=h4_bold,
                h5_bold=h5_bold,
                line_spacing=line_spacing,
                space_after_pt=space_after_pt,
                font_name=font_name,
                font_color=font_color,
            )


def _add_content_line(
    text_frame: "TextFrame",
    line: str,
    *,
    h2_size: int,
    h3_size: int,
    h4_size: int,
    h5_size: int,
    body_size: int,
    bullet_size: int,
    numbered_size: int,
    spacer_size: int,
    numbering_type: str,
    h2_bold: bool,
    h3_bold: bool,
    h4_bold: bool,
    h5_bold: bool,
    line_spacing: float | None = None,
    space_after_pt: int | None = None,
    font_name: str | None = None,
    font_color: RGBColor | None = None,
) -> None:
    """Add a single content line to a text frame with appropriate formatting.
    
    Internal helper function that handles the different line types:
    spacers, headers, bullets, numbered lists, and plain text.
    
    Args:
        text_frame: PowerPoint text frame to add content to.
        line: Single stripped line of content.
        h2_size - h5_bold: Formatting parameters from config.
        line_spacing: Line spacing multiplier (1.0 = single, None = template default).
        space_after_pt: Space after paragraph in points (None = template default).
        font_name: Font family name (e.g. "Arial"). Applied to all runs as a
            fallback from the template when config doesn't override.
        font_color: Font color as :class:`~pptx.util.RGBColor`. Applied to all
            runs as a fallback from the template when config doesn't override.
    """
    def _apply_paragraph_spacing(para):
        """Apply line spacing and space_after to a paragraph if configured."""
        if line_spacing is not None:
            para.line_spacing = line_spacing
        if space_after_pt is not None:
            para.space_after = Pt(space_after_pt)
    
    def _apply_template_font(run) -> None:
        """Apply template-derived font defaults (name, color) to a run.

        Skips font-name override when the run already carries an explicit
        font (e.g. 'Consolas' set by inline-code formatting in rich_text).
        """
        if font_name is not None and run.font.name is None:
            run.font.name = font_name
        if font_color is not None:
            run.font.color.rgb = font_color
    
    # Handle spacer markers (blank lines in markdown)
    if line == SPACER_MARKER:
        p = text_frame.add_paragraph()
        # Use a space character to ensure the paragraph renders with height
        p.text = " "
        remove_bullet(p)
        for run in p.runs:
            run.font.size = Pt(spacer_size)
        p.space_before = Pt(spacer_size)
        p.space_after = Pt(0)
        logger.debug(f"  Added spacer paragraph ({spacer_size}pt)")
        return
    
    # Handle H5 header (##### )
    if line.startswith("##### "):
        p = text_frame.add_paragraph()
        add_formatted_text(p, line[6:])
        p.level = 0
        remove_bullet(p)
        _apply_paragraph_spacing(p)
        for run in p.runs:
            run.font.size = Pt(h5_size)
            if h5_bold:
                run.font.bold = True
            _apply_template_font(run)
        logger.debug(f"  Added H5: {line[6:]}")
        return
    
    # Handle H4 header (#### )
    if line.startswith("#### "):
        p = text_frame.add_paragraph()
        add_formatted_text(p, line[5:])
        p.level = 0
        remove_bullet(p)
        _apply_paragraph_spacing(p)
        for run in p.runs:
            run.font.size = Pt(h4_size)
            if h4_bold:
                run.font.bold = True
            _apply_template_font(run)
        logger.debug(f"  Added H4: {line[5:]}")
        return
    
    # Handle H3 header (### )
    if line.startswith("### "):
        p = text_frame.add_paragraph()
        add_formatted_text(p, line[4:])
        p.level = 0
        remove_bullet(p)
        _apply_paragraph_spacing(p)
        for run in p.runs:
            run.font.size = Pt(h3_size)
            if h3_bold:
                run.font.bold = True
            _apply_template_font(run)
        logger.debug(f"  Added H3: {line[4:]}")
        return
    
    # Handle H2 header (## )
    if line.startswith("## "):
        p = text_frame.add_paragraph()
        add_formatted_text(p, line[3:])
        p.level = 0
        remove_bullet(p)
        _apply_paragraph_spacing(p)
        for run in p.runs:
            run.font.size = Pt(h2_size)
            if h2_bold:
                run.font.bold = True
            _apply_template_font(run)
        logger.debug(f"  Added H2: {line[3:]}")
        return
    
    # Handle level-1 bullet (indented: "  - ")
    if line.startswith("  - "):
        p = text_frame.add_paragraph()
        add_formatted_text(p, line[4:])
        p.level = 1
        add_bullet(p, level=1)
        _apply_paragraph_spacing(p)
        for run in p.runs:
            run.font.size = Pt(bullet_size)
            _apply_template_font(run)
        logger.debug(f"  Added sub-bullet: {line[4:]}")
        return
    
    # Handle level-0 bullet ("- ")
    if line.startswith("- "):
        p = text_frame.add_paragraph()
        add_formatted_text(p, line[2:])
        p.level = 0
        add_bullet(p, level=0)
        _apply_paragraph_spacing(p)
        for run in p.runs:
            run.font.size = Pt(bullet_size)
            _apply_template_font(run)
        logger.debug(f"  Added bullet: {line[2:]}")
        return
    
    # Handle numbered lists (e.g., "1. ", "2. ")
    numbered_match = re.match(r"^(\d+)\.\s+(.*)$", line)
    if numbered_match:
        num = int(numbered_match.group(1))
        text = numbered_match.group(2)
        p = text_frame.add_paragraph()
        add_formatted_text(p, text)
        p.level = 0
        add_numbering(p, start_at=num, numbering_type=numbering_type)
        _apply_paragraph_spacing(p)
        for run in p.runs:
            run.font.size = Pt(numbered_size)
            _apply_template_font(run)
        logger.debug(f"  Added numbered item {num}: {text}")
        return
    
    # Plain text (default)
    p = text_frame.add_paragraph()
    add_formatted_text(p, line)
    remove_bullet(p)
    _apply_paragraph_spacing(p)
    for run in p.runs:
        run.font.size = Pt(body_size)
        _apply_template_font(run)
    logger.debug(f"  Added text: {line}")
